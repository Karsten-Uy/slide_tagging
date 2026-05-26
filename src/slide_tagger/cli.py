"""Command-line entry point for Pipeline A.

    slide-tagger tag deck.pptx              # paste-ready STRUCTURAL DATA blocks
    slide-tagger tag deck.pptx --slide 2    # just one slide
    slide-tagger tag deck.pptx --json       # full structural JSON
    slide-tagger deck-summary deck.pptx     # paste-ready DECK SUMMARY block
    slide-tagger template deck.pptx         # blank hand-tagging template (JSON)
    slide-tagger validate labels.json       # validate a hand-tagged file
    slide-tagger render deck.pptx           # per-slide PNGs (full + thumbnail)
    slide-tagger extract-assets deck.pptx   # extract recurring logos/branding images -> PNGs
    slide-tagger merge vlm.json input.json  # re-impose Pipeline A fields on VLM output
    slide-tagger score pred.json truth.json # score enriched output vs hand-label
    slide-tagger eval                       # score data/tagged/* vs hand_labels/*
    slide-tagger bench --runs 3             # API: N runs/deck, mean±std (needs ANTHROPIC_API_KEY)

`tag`/`deck-summary` emit paste-ready `STRUCTURAL DATA` / `DECK SUMMARY` grounding
blocks for inspecting Pipeline A's output. `template`/`validate` drive the
enrichment + hand-labeling workflow (docs/vlm_prompt_test.md, docs/manual_tagging.md,
docs/deck_tagging_prompt.md) — `template` emits a blank record (structural filled,
enrichment null) and `validate` checks it against the schema (init.md Phase 1).
"""

from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

from pydantic import ValidationError

from slide_tagger.eval import score_corpus, score_deck
from slide_tagger.eval.report import render_console, render_json, render_markdown
from slide_tagger.extractors.render.paths import deck_slug, render_rel_path, thumb_rel_path
from slide_tagger.extractors.structural.aggregator import summarize_deck
from slide_tagger.extractors.structural.pptx_parser import parse_pptx
from slide_tagger.merge import merge_structural
from slide_tagger.enrich import enrich_once, upload_pdf
from slide_tagger.prompt_source import resolve_prompt
from slide_tagger.provenance import build_enriched_record, is_filled
from slide_tagger.schema.enums import DensityBucket
from slide_tagger.schema.models import DeckStructural, DeckSummary, SlideStructural
from slide_tagger.schema.tagged import DeckTag, blank_tag, legend


def structural_data_block(slide: SlideStructural) -> str:
    """Format one slide as the paste-ready STRUCTURAL DATA block."""
    d = slide.density
    title = slide.title_text or ""
    title_pos = slide.title_position.value if slide.title_position else ""
    charts = "yes" if slide.has_chart else "no"
    tables = "yes" if slide.has_table else "no"
    return "\n".join(
        [
            "STRUCTURAL DATA (from Pipeline A — context for your judgments; do NOT output these fields):",
            f"  word_count: {d.word_count}",
            f"  text_blocks: {d.text_blocks}",
            f"  visual_elements: {d.visual_elements}",
            f"  whitespace_ratio_est: {d.whitespace_ratio_est}",
            f"  density_bucket: {d.bucket.value}",
            f'  title_text: "{title}"',
            f"  title_position: {title_pos}",
            f"  images: {slide.image_count}   charts: {charts}   tables: {tables}",
        ]
    )


def deck_summary_block(summary: DeckSummary) -> str:
    """Format a deck as the paste-ready DECK SUMMARY block (deck-level pass)."""
    dd = summary.density_distribution
    dist = ", ".join(f"{b.value}={dd.get(b, 0)}" for b in DensityBucket)
    pos = summary.dominant_title_position.value if summary.dominant_title_position else ""
    return "\n".join(
        [
            "DECK SUMMARY (from Pipeline A — deterministic context; do NOT output these fields):",
            f"  slide_count: {summary.slide_count}",
            f"  density_distribution: {dist}",
            f"  slides_with_charts: {summary.slides_with_charts}   "
            f"slides_with_tables: {summary.slides_with_tables}   "
            f"slides_with_images: {summary.slides_with_images}",
            f"  dominant_title_position: {pos}",
            f"  avg_word_count: {summary.avg_word_count}",
        ]
    )


def _resolve_deck(deck: Path) -> Path | int:
    """Validate a deck path. Returns the path, or an error exit code."""
    suffix = deck.suffix.lower()
    if suffix == ".pdf":
        print(
            "PDF parsing is not implemented yet — Pipeline A currently supports .pptx. "
            "(See init.md open question on PDF support priority.)",
            file=sys.stderr,
        )
        return 2
    if suffix != ".pptx":
        print(f"Unsupported file type '{suffix}'. Expected a .pptx file.", file=sys.stderr)
        return 2
    if not deck.exists():
        print(f"File not found: {deck}", file=sys.stderr)
        return 2
    return deck


def _cmd_tag(args: argparse.Namespace) -> int:
    resolved = _resolve_deck(args.deck)
    if isinstance(resolved, int):
        return resolved
    deck = resolved

    result = parse_pptx(deck)

    if args.json:
        print(result.model_dump_json(indent=2))
        return 0

    slides = result.slides
    if args.slide is not None:
        if not 0 <= args.slide < len(slides):
            print(
                f"Slide index {args.slide} out of range (0..{len(slides) - 1}).",
                file=sys.stderr,
            )
            return 2
        slides = [slides[args.slide]]

    for slide in slides:
        print(f"# slide {slide.index}  ({result.source_filename})")
        print(structural_data_block(slide))
        print()
    return 0


def _cmd_deck_summary(args: argparse.Namespace) -> int:
    resolved = _resolve_deck(args.deck)
    if isinstance(resolved, int):
        return resolved

    summary = summarize_deck(parse_pptx(resolved))

    if args.json:
        print(summary.model_dump_json(indent=2))
        return 0

    print(deck_summary_block(summary))
    return 0


def _load_structural(src: Path) -> DeckStructural | int:
    """Load structural data from a .pptx (parse) or a Pipeline A .json (load)."""
    suffix = src.suffix.lower()
    if suffix == ".json":
        if not src.exists():
            print(f"File not found: {src}", file=sys.stderr)
            return 2
        try:
            return DeckStructural.model_validate_json(src.read_text(encoding="utf-8"))
        except ValidationError as exc:
            print(f"Not a valid Pipeline A structural JSON:\n{exc}", file=sys.stderr)
            return 1
    resolved = _resolve_deck(src)
    if isinstance(resolved, int):
        return resolved
    return parse_pptx(resolved)


def _attach_render_paths(tag: DeckTag) -> None:
    """Fill each slide's render_path/thumbnail_path from the deck slug + index, so
    the template already references where renders live (whether or not they exist
    yet). The downstream MCP server resolves them via THUMBNAIL_BASE_PATH."""
    slug = deck_slug(tag.source_filename)
    for slide in tag.slides:
        slide.render_path = render_rel_path(slug, slide.index)
        slide.thumbnail_path = thumb_rel_path(slug, slide.index)


def _cmd_template(args: argparse.Namespace) -> int:
    deck = _load_structural(args.source)
    if isinstance(deck, int):
        return deck
    tag = blank_tag(deck)
    _attach_render_paths(tag)
    out = {"_legend": legend(), **tag.model_dump(mode="json")}
    print(json.dumps(out, indent=2, ensure_ascii=False))
    return 0


def _cmd_render(args: argparse.Namespace) -> int:
    resolved = _resolve_deck(args.deck)
    if isinstance(resolved, int):
        return resolved

    # Imported lazily so `tag`/`template` work even if pdf2image isn't installed.
    from slide_tagger.extractors.render import render_deck
    from slide_tagger.extractors.render.soffice import LibreOfficeNotFound, RenderError

    try:
        result = render_deck(
            resolved,
            out_root=args.out,
            dpi=args.dpi,
            thumb_px=args.thumb,
            only_index=args.slide,
            poppler_path=args.poppler_path,
        )
    except LibreOfficeNotFound as exc:
        print(str(exc), file=sys.stderr)
        return 2
    except RenderError as exc:
        print(f"Render failed: {exc}", file=sys.stderr)
        return 1
    except Exception as exc:  # pdf2image / poppler failures surface here
        print(
            f"Render failed (is poppler installed? pass --poppler-path to its bin/): {exc}",
            file=sys.stderr,
        )
        return 1

    dest = Path(args.out) / result.deck_slug
    print(f"# rendered {len(result.slides)} slide(s) from {Path(resolved).name} -> {dest}")
    for slide in result.slides:
        print(f"  slide {slide.index}: {slide.render_path}  ·  {slide.thumbnail_path}")
    return 0


def _write_contact_sheet(out_dir: Path, elements: list[dict]) -> None:
    """Write a captioned grid of the extracted images to help hand-labeling."""
    from PIL import Image, ImageDraw

    tiles = []
    for e in elements:
        p = out_dir / Path(e["image_path"]).name
        try:
            tiles.append((e, Image.open(p).convert("RGB")))
        except Exception:
            continue
    if not tiles:
        return
    cw, cap, cols = 240, 48, 3
    rows = (len(tiles) + cols - 1) // cols
    ch = cw + cap
    sheet = Image.new("RGB", (cw * cols, ch * rows), "white")
    draw = ImageDraw.Draw(sheet)
    for n, (e, im) in enumerate(tiles):
        im.thumbnail((cw - 10, cw - 10))
        r, c = divmod(n, cols)
        x, y = c * cw, r * ch
        sheet.paste(im, (x + 5, y + 5))
        cap_txt = f"#{n} {e.get('type') or '?'} ({e.get('source')}, {len(e.get('appears_on_slides') or [])} sl)"
        draw.text((x + 5, y + cw - 8), cap_txt[:42], fill="#222")
    sheet.save(out_dir / "_contactsheet.png")


def _merge_recurring_into(path: Path, elements: list[dict]) -> int:
    """Append image-based recurring elements into a hand-label JSON (dedupe by phash)."""
    if not path.exists():
        print(f"--into file not found: {path}", file=sys.stderr)
        return 2
    data = json.loads(path.read_text(encoding="utf-8"))
    existing = data.setdefault("design_system", {}).setdefault("recurring_elements", [])
    have = {r.get("phash") for r in existing if r.get("phash")}
    added = 0
    for e in elements:
        if e.get("phash") and e["phash"] in have:
            continue
        existing.append(e)
        added += 1
    path.write_text(json.dumps(data, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")
    print(f"# merged {added} image element(s) into {path.name}", file=sys.stderr)
    return 0


def _cmd_extract_assets(args: argparse.Namespace) -> int:
    from io import BytesIO

    from PIL import Image
    from pptx import Presentation

    from slide_tagger.extractors.structural.recurring_images import extract_recurring_images

    resolved = _resolve_deck(args.deck)
    if isinstance(resolved, int):
        return resolved

    prs = Presentation(str(resolved))
    slug = deck_slug(Path(resolved).name)
    groups = extract_recurring_images(prs, fraction=args.fraction)

    out_dir = args.out / slug
    out_dir.mkdir(parents=True, exist_ok=True)
    elements: list[dict] = []
    for n, g in enumerate(groups):
        fname = f"recurring_{n:02d}.png"
        try:
            with Image.open(BytesIO(g.blob)) as im:
                im.convert("RGB").save(out_dir / fname)
        except Exception as exc:  # vector/unsupported → manual fallback
            print(f"# skip group {n}: cannot rasterize ({exc})", file=sys.stderr)
            continue
        rel = f"assets/{slug}/{fname}"
        elements.append(
            {
                "type": g.type.value if g.type else None,
                "value": None,
                "phash": g.phash,
                "position": g.position.value if g.position else None,
                "appears_on_slides": g.slide_indices,
                "image_path": rel,
                "source": g.source,
            }
        )
        print(
            f"# saved {rel}  type={g.type.value if g.type else 'None'}  "
            f"source={g.source}  slides={len(g.slide_indices)}",
            file=sys.stderr,
        )

    if not elements:
        print(
            f"# {slug}: no raster branding images extracted — lower --fraction or use "
            "the manual fallback (docs/manual_tagging.md).",
            file=sys.stderr,
        )
    if args.contact_sheet and elements:
        _write_contact_sheet(out_dir, elements)
        print(f"# wrote {out_dir / '_contactsheet.png'}", file=sys.stderr)

    print(json.dumps(elements, indent=2, ensure_ascii=False))
    if args.into is not None and elements:
        return _merge_recurring_into(args.into, elements)
    return 0


# `is_filled` moved to provenance.py (shared with the paste harness); aliased for
# any callers/tests still expecting the private name.
_is_filled = is_filled


def _load_tag(path: Path) -> DeckTag | int:
    """Load a tagged-deck JSON (stripping the template `_legend`) into a DeckTag,
    or return an error exit code (2 not-found, 1 invalid JSON / schema)."""
    if not path.exists():
        print(f"File not found: {path}", file=sys.stderr)
        return 2
    try:
        data = json.loads(path.read_text(encoding="utf-8"))
    except json.JSONDecodeError as exc:
        print(f"Invalid JSON in {path.name}: {exc}", file=sys.stderr)
        return 1
    data.pop("_legend", None)  # template helper, not part of the schema
    try:
        return DeckTag.model_validate(data)
    except ValidationError as exc:
        print(f"✗ Schema errors in {path.name}:\n{exc}", file=sys.stderr)
        return 1


def _cmd_validate(args: argparse.Namespace) -> int:
    path: Path = args.labels
    tag = _load_tag(path)
    if isinstance(tag, int):
        return tag

    # Deck-level enrichment fields the tagger fills in (deck_length is auto-set
    # from slide_count, so it isn't counted here).
    deck_fields = {
        "client_industry": tag.client_industry,
        "client_sub_industry": tag.client_sub_industry,
        "client_type": tag.client_type,
        "engagement_stage": tag.engagement_stage,
        "content_area": tag.content_area,
        "audience_level": tag.audience_level,
        "deliverable_format": tag.deliverable_format,
        "geography": tag.geography,
        "confidentiality_tier": tag.confidentiality_tier,
        "inferred_publisher": tag.inferred_publisher,
        "deck_summary_one_sentence": tag.deck_summary_one_sentence,
    }
    deck_missing = [k for k, v in deck_fields.items() if not _is_filled(v)]

    # A slide counts as fully tagged once its core enrichment fields are set.
    incomplete = [
        s.index
        for s in tag.slides
        if not _is_filled(s.slide_purpose)
        or not _is_filled(s.message_type)
        or not _is_filled(s.main_message)
        or not _is_filled(s.dominant_visual_element)
    ]

    print(f"✓ Schema valid: {path.name}")
    deck_done = len(deck_fields) - len(deck_missing)
    line = f"Deck-level: {deck_done}/{len(deck_fields)} filled"
    if deck_missing:
        line += f"  (missing: {', '.join(deck_missing)})"
    print(line)
    done = len(tag.slides) - len(incomplete)
    line = f"Slides fully tagged: {done}/{len(tag.slides)}"
    if incomplete:
        line += f"  (incomplete: {incomplete})"
    print(line)

    ds = tag.design_system
    if ds is not None and ds.recurring_elements:
        untyped = [i for i, r in enumerate(ds.recurring_elements) if r.type is None]
        typed = len(ds.recurring_elements) - len(untyped)
        line = f"Recurring elements typed: {typed}/{len(ds.recurring_elements)}"
        if untyped:
            line += f"  (untyped indices: {untyped})"
        if ds.grid is None:
            line += "  · grid: unset"
        print(line)

    ir = tag.inferred_rules
    ir_populated = ir is not None and (
        _is_filled(ir.title.font_family_observed)
        or _is_filled(ir.title.size_pt_most_common)
        or _is_filled(ir.body_text.font_family_observed)
        or _is_filled(ir.color_palette.primary_observed)
    )
    prov_set = tag.provenance is not None and _is_filled(tag.provenance.tagged_by)
    print(
        f"Inferred rules: {'populated' if ir_populated else 'empty'}  ·  "
        f"provenance: {'set' if prov_set else 'unset'}"
    )
    return 0


def _read_json_lenient(path: Path) -> dict | int:
    """Read a JSON file into a raw dict, tolerating a leading ```json code fence
    (the web GUI often wraps output in one) and stripping the template `_legend`.
    Returns an error exit code on failure."""
    if not path.exists():
        print(f"File not found: {path}", file=sys.stderr)
        return 2
    text = path.read_text(encoding="utf-8").strip()
    if text.startswith("```"):
        lines = text.splitlines()
        lines = lines[1:]  # drop the opening ``` / ```json
        if lines and lines[-1].strip().startswith("```"):
            lines = lines[:-1]
        text = "\n".join(lines).strip()
    try:
        data = json.loads(text)
    except json.JSONDecodeError as exc:
        print(f"Invalid JSON in {path.name}: {exc}", file=sys.stderr)
        return 1
    if not isinstance(data, dict):
        print(f"Expected a JSON object in {path.name}, got {type(data).__name__}.", file=sys.stderr)
        return 1
    data.pop("_legend", None)
    return data


def _cmd_merge(args: argparse.Namespace) -> int:
    vlm = _read_json_lenient(args.vlm_output)
    if isinstance(vlm, int):
        return vlm
    template = _read_json_lenient(args.template)
    if isinstance(template, int):
        return template

    merged = merge_structural(vlm, template)

    out_text = json.dumps(merged, indent=2, ensure_ascii=False)
    if args.out is not None:
        args.out.parent.mkdir(parents=True, exist_ok=True)
        args.out.write_text(out_text + "\n", encoding="utf-8")
        print(f"# wrote merged record -> {args.out}", file=sys.stderr)
    else:
        print(out_text)

    # Best-effort validation so the user sees any *remaining* (non-structural)
    # problems — e.g. an out-of-enum value the VLM produced — right away.
    try:
        DeckTag.model_validate(merged)
        print("# structural fields re-imposed from template; result validates.", file=sys.stderr)
    except ValidationError as exc:
        n = len(exc.errors())
        print(
            f"# structural fields re-imposed, but {n} schema issue(s) remain "
            "(not structural — fix in the VLM output or prompt):",
            file=sys.stderr,
        )
        for err in exc.errors()[:10]:
            print(f"#   {'.'.join(str(x) for x in err['loc'])}: {err['msg']}", file=sys.stderr)
    return 0


# Maps a hand-label stem (reference_data/hand_labels/<stem>.tagged.json) to the
# source file stem in reference_data/hand_labels/sources/ (they differ for
# ereadiness). The clean, index-aligned decks default on; the others are opt-in.
_BENCH_DECKS = {
    "nigeria-economic-outlook-october-2023-v1": "nigeria-economic-outlook-october-2023-v1",
    "digital-auto-report-2023": "digital-auto-report-2023",
    "electric-vehicle-sales-review-q4-2022": "electric-vehicle-sales-review-q4-2022",
    "ereadiness-study-2023": "strategyand-ereadiness-study-2023",
}
_BENCH_DEFAULT = ["nigeria-economic-outlook-october-2023-v1", "digital-auto-report-2023"]


def _build_template(pptx: Path) -> dict:
    """Pipeline A structural template + _legend, identical to `template` output —
    but with the 6 §A pre-fillable fields populated per slide. The merge guard
    (see `merge._SLIDE_STRUCTURAL`) will re-impose these after the VLM round-trip
    so a VLM that 'helpfully' overwrote a pre-fill can't corrupt them."""
    from slide_tagger.extractors.structural.prefill import prefill_deck

    tag = blank_tag(parse_pptx(pptx))
    _attach_render_paths(tag)
    out = {"_legend": legend(), **tag.model_dump(mode="json")}

    prefilled = prefill_deck(pptx)
    out_slides = out.get("slides", [])
    for slide_dict, fields in zip(out_slides, prefilled):
        for k, v in fields.items():
            slide_dict[k] = v
    return out


def _cmd_bench(args: argparse.Namespace) -> int:
    import statistics
    import time
    from collections import defaultdict

    try:
        import anthropic
    except ImportError:
        print("`bench` needs the anthropic SDK — run `uv sync` (it's now a dependency).", file=sys.stderr)
        return 2

    try:
        system = resolve_prompt(args.prompt).text
    except FileNotFoundError as exc:
        print(str(exc), file=sys.stderr)
        return 2

    try:
        # Hard per-request timeout so a network hang errors out instead of
        # stalling forever (the first uncached call is the slow one).
        client = anthropic.Anthropic(timeout=float(args.timeout))
    except Exception as exc:  # missing/invalid key surfaces here
        print(f"Could not init Anthropic client (is ANTHROPIC_API_KEY set?): {exc}", file=sys.stderr)
        return 2

    decks = args.deck or _BENCH_DEFAULT
    deck_scores: dict[str, list[float]] = {}
    per_field: dict[str, list[float]] = defaultdict(list)

    for stem in decks:
        src = _BENCH_DECKS.get(stem, stem)
        pptx, pdf = args.sources / f"{src}.pptx", args.sources / f"{src}.pdf"
        label = args.labels / f"{stem}.tagged.json"
        if not (pptx.exists() and pdf.exists() and label.exists()):
            print(f"# skip {stem}: need {pptx.name}, {pdf.name}, {label.name}", file=sys.stderr)
            continue
        truth = _load_tag(label)
        if isinstance(truth, int):
            return truth
        template = _build_template(pptx)
        template_core = {k: v for k, v in template.items() if k != "_legend"}

        print(f"# {stem}: uploading PDF ({pdf.stat().st_size // 1024} KB)…", file=sys.stderr)
        try:
            file_id = upload_pdf(client, pdf)
        except Exception as exc:  # auth / network / API errors surface here
            print(f"# {stem}: upload failed ({exc}). Is ANTHROPIC_API_KEY set and valid?", file=sys.stderr)
            return 2
        print(f"# {stem}: uploaded (file_id={file_id}).", file=sys.stderr)
        (args.out / stem).mkdir(parents=True, exist_ok=True)
        scores: list[float] = []
        for k in range(args.runs):
            print(
                f"# {stem}: run {k + 1}/{args.runs} — calling {args.model} "
                f"(effort={args.effort}); '·'=thinking, '.'=output…",
                file=sys.stderr,
            )
            t0 = time.perf_counter()
            try:
                enriched = enrich_once(
                    client, system=system, template=template,
                    file_id=file_id, model=args.model, effort=args.effort,
                    verbose=not args.quiet,
                )
            except Exception as exc:
                print(f"# {stem} run {k + 1}/{args.runs} FAILED after {time.perf_counter() - t0:.0f}s: {exc}", file=sys.stderr)
                continue
            elapsed = time.perf_counter() - t0
            enriched.pop("_legend", None)
            merged = merge_structural(enriched, template_core)
            try:
                pred = DeckTag.model_validate(merged)
            except ValidationError as exc:
                print(f"# {stem} run {k + 1}: invalid after sanitize ({len(exc.errors())} errs) — skipped", file=sys.stderr)
                continue
            ds = score_deck(pred, truth, name=f"{stem}#{k + 1}")
            acc = score_corpus([ds]).headline_accuracy or 0.0
            scores.append(acc)
            for path, fr in ds.results.items():
                if fr.scored and fr.accuracy is not None:
                    per_field[path].append(fr.accuracy)
            (args.out / stem / f"run_{k + 1}.json").write_text(
                json.dumps(merged, indent=2, ensure_ascii=False), encoding="utf-8"
            )
            print(f"# {stem} run {k + 1}/{args.runs}: {acc:.0%}  ({elapsed:.0f}s)", file=sys.stderr)
        try:
            client.beta.files.delete(file_id)
        except Exception:
            pass
        if scores:
            deck_scores[stem] = scores

    if not deck_scores:
        print("No successful runs.", file=sys.stderr)
        return 1

    print("\n=== bench report ===")
    print(f"prompt: {args.prompt}  ·  model: {args.model}  ·  effort: {args.effort}  ·  runs/deck: {args.runs}")
    deck_means = []
    for stem, scores in deck_scores.items():
        mean = statistics.mean(scores)
        sd = statistics.pstdev(scores) if len(scores) > 1 else 0.0
        deck_means.append(mean)
        runs_str = ", ".join(f"{s:.0%}" for s in scores)
        print(f"  {stem}: {mean:.1%} ± {sd:.1%}  (runs: {runs_str})")
    if deck_means:
        print(f"  CORPUS (mean of deck means): {statistics.mean(deck_means):.1%}")

    field_means = {p: statistics.mean(v) for p, v in per_field.items()}
    print("\nPer-field mean accuracy across all runs (weakest first):")
    for path, mean in sorted(field_means.items(), key=lambda kv: kv[1]):
        print(f"  {path:48} {mean:.0%}  (n_runs={len(per_field[path])})")

    # --- persist: optional JSON + an always-on results-log row ---
    import datetime

    corpus = statistics.mean(deck_means)
    deck_summary = {
        stem: {
            "mean": statistics.mean(s),
            "std": statistics.pstdev(s) if len(s) > 1 else 0.0,
            "runs": s,
        }
        for stem, s in deck_scores.items()
    }
    stamp = datetime.datetime.now().isoformat(timespec="seconds")
    if args.json:
        args.json.parent.mkdir(parents=True, exist_ok=True)
        args.json.write_text(
            json.dumps(
                {"label": args.label, "timestamp": stamp, "model": args.model,
                 "effort": args.effort, "runs": args.runs, "corpus": corpus,
                 "decks": deck_summary, "per_field": field_means},
                indent=2,
            ),
            encoding="utf-8",
        )
        print(f"\n# wrote JSON report -> {args.json}", file=sys.stderr)

    weakest = "; ".join(f"{p} {m:.0%}" for p, m in sorted(field_means.items(), key=lambda kv: kv[1])[:3])
    decks_cell = "; ".join(
        f"{stem.split('-')[0]} {d['mean']:.1%}±{d['std']:.1%}" for stem, d in deck_summary.items()
    )
    row = f"| {args.label or '—'} | {stamp[:10]} | {args.model} | {args.effort} | {args.runs} | {decks_cell} | {corpus:.1%} | {weakest} |\n"
    if not args.log.exists():
        args.log.parent.mkdir(parents=True, exist_ok=True)
        args.log.write_text(
            "# Bench results log\n\n"
            "| label | date | model | effort | runs | per-deck mean±std | corpus | weakest 3 fields |\n"
            "|---|---|---|---|---|---|---|---|\n",
            encoding="utf-8",
        )
    with open(args.log, "a", encoding="utf-8") as f:
        f.write(row)
    print(f"# appended results row -> {args.log}", file=sys.stderr)
    return 0


# --- enrich (automated Pipeline A+B tagger) ---------------------------------

# Enrichment-fields constants + record assembly moved to `provenance.py` (shared
# with the paste harness so both produce identical records). The aliases here
# preserve any external callers/tests that imported the private names.
from slide_tagger.provenance import (
    DECK_ENRICHMENT_FIELDS as _DECK_ENRICHMENT,
    SLIDE_ENRICHMENT_FIELDS as _SLIDE_ENRICHMENT,
    change_field as _change_field,
    filled_enrichment_fields as _filled_enrichment_fields,
    unfilled_enrichment_fields as _unfilled_enrichment_fields,
)

_build_enriched_record = build_enriched_record


def _finalize_enrich(
    enriched: dict,
    template_core: dict,
    sanitizer_changes: list[str],
    artifact,  # PromptArtifact
    model: str,
    tagged_by: str | None = None,
) -> DeckTag:
    """Build the stamped record and validate it (raises ValidationError on failure).
    Convenience wrapper around `build_enriched_record` for callers/tests that want a
    validated `DeckTag`."""
    return DeckTag.model_validate(
        build_enriched_record(enriched, template_core, sanitizer_changes, artifact, model, tagged_by)
    )


def _convert_pptx_to_pdf(pptx: Path, out_dir: Path) -> Path:
    """Thin wrapper around the LibreOffice pptx→PDF conversion — a clean monkeypatch
    point for tests. Raises LibreOfficeNotFound / RenderError from the soffice module."""
    from slide_tagger.extractors.render.soffice import pptx_to_pdf

    return pptx_to_pdf(pptx, out_dir)


def _write_into_corpus(pptx: Path, record: dict, corpus_dir: Path, *, fraction: float) -> None:
    """Copy the tagged record into the served corpus dir and extract logo PNGs into
    <corpus>/assets/<slug>/, merging the image elements (deduped by phash) into the
    corpus copy. Never touches reference_data/hand_labels (the eval answer key)."""
    import copy
    from io import BytesIO

    from PIL import Image
    from pptx import Presentation

    from slide_tagger.extractors.structural.recurring_images import extract_recurring_images

    slug = deck_slug(pptx.name)
    corpus_dir.mkdir(parents=True, exist_ok=True)

    groups = extract_recurring_images(Presentation(str(pptx)), fraction=fraction)
    assets_dir = corpus_dir / "assets" / slug
    if groups:
        assets_dir.mkdir(parents=True, exist_ok=True)
    new_elems: list[dict] = []
    for n, g in enumerate(groups):
        fname = f"recurring_{n:02d}.png"
        try:
            with Image.open(BytesIO(g.blob)) as im:
                im.convert("RGB").save(assets_dir / fname)
        except Exception as exc:  # vector/unsupported → manual fallback
            print(f"# skip group {n}: cannot rasterize ({exc})", file=sys.stderr)
            continue
        new_elems.append(
            {
                "type": g.type.value if g.type else None,
                "value": None,
                "phash": g.phash,
                "position": g.position.value if g.position else None,
                "appears_on_slides": g.slide_indices,
                "image_path": f"assets/{slug}/{fname}",
                "source": g.source,
            }
        )

    corpus_record = copy.deepcopy(record)
    ds = corpus_record.get("design_system") or {}
    corpus_record["design_system"] = ds
    existing = ds.setdefault("recurring_elements", [])
    have = {r.get("phash") for r in existing if r.get("phash")}
    added = 0
    for e in new_elems:
        if e.get("phash") and e["phash"] in have:
            continue
        existing.append(e)
        added += 1

    dest = corpus_dir / f"{slug}.tagged.json"
    dest.write_text(
        json.dumps({"_legend": legend(), **corpus_record}, indent=2, ensure_ascii=False) + "\n",
        encoding="utf-8",
    )
    print(f"# corpus: wrote {dest}  (+{added} logo asset(s))", file=sys.stderr)


def _cmd_enrich(args: argparse.Namespace) -> int:
    import os
    import tempfile

    resolved = _resolve_deck(args.deck)
    if isinstance(resolved, int):
        return resolved
    pptx = resolved

    try:
        import anthropic
    except ImportError:
        print("`enrich` needs the anthropic SDK — run `uv sync` (it's a dependency).", file=sys.stderr)
        return 2

    try:
        artifact = resolve_prompt(args.prompt)
    except FileNotFoundError as exc:
        print(str(exc), file=sys.stderr)
        return 2

    try:
        client = anthropic.Anthropic(timeout=float(args.timeout))
    except Exception as exc:  # missing/invalid key surfaces here
        print(f"Could not init Anthropic client (is ANTHROPIC_API_KEY set?): {exc}", file=sys.stderr)
        return 2

    template = _build_template(pptx)
    template_core = {k: v for k, v in template.items() if k != "_legend"}

    from slide_tagger.extractors.render.soffice import LibreOfficeNotFound, RenderError

    def _do_enrich(pdf_path: Path):
        print(f"# uploading PDF ({pdf_path.stat().st_size // 1024} KB)…", file=sys.stderr)
        file_id = upload_pdf(client, pdf_path)
        print(
            f"# uploaded (file_id={file_id}); calling {args.model} (effort={args.effort})…",
            file=sys.stderr,
        )
        try:
            return enrich_once(
                client, system=artifact.text, template=template, file_id=file_id,
                model=args.model, effort=args.effort, verbose=not args.quiet,
                return_changes=True,
            )
        finally:
            try:
                client.beta.files.delete(file_id)
            except Exception:
                pass

    try:
        if args.pdf is not None:
            if not args.pdf.exists():
                print(f"--pdf not found: {args.pdf}", file=sys.stderr)
                return 2
            enriched, changes = _do_enrich(args.pdf)
        else:
            with tempfile.TemporaryDirectory(prefix="enrich_pdf_") as tmp:
                try:
                    pdf_path = _convert_pptx_to_pdf(pptx, Path(tmp))
                except LibreOfficeNotFound as exc:
                    print(f"{exc}\nInstall LibreOffice or pass --pdf <file>.", file=sys.stderr)
                    return 2
                except RenderError as exc:
                    print(f"pptx→PDF conversion failed: {exc}", file=sys.stderr)
                    return 1
                enriched, changes = _do_enrich(pdf_path)
    except Exception as exc:  # upload / API errors surface here
        print(f"# enrich failed ({exc}). Is ANTHROPIC_API_KEY set and valid?", file=sys.stderr)
        return 1

    record = _build_enriched_record(
        enriched, template_core, changes, artifact, args.model, tagged_by=args.tagged_by
    )
    try:
        DeckTag.model_validate(record)
    except ValidationError as exc:
        print(f"# {len(exc.errors())} schema issue(s) after merge (not structural):", file=sys.stderr)
        for err in exc.errors()[:10]:
            print(f"#   {'.'.join(str(x) for x in err['loc'])}: {err['msg']}", file=sys.stderr)
        if args.strict:
            return 1

    out_path = args.out or (Path("data/tagged") / f"{deck_slug(pptx.name)}.tagged.json")
    out_path.parent.mkdir(parents=True, exist_ok=True)
    out_path.write_text(
        json.dumps({"_legend": legend(), **record}, indent=2, ensure_ascii=False) + "\n",
        encoding="utf-8",
    )
    n_low = len(record["provenance"]["low_confidence_fields"])
    print(f"# wrote tagged record -> {out_path}", file=sys.stderr)
    print(
        f"# prompt_version={artifact.version}  model={args.model}  low_confidence_fields={n_low}",
        file=sys.stderr,
    )

    if args.into_corpus:
        corpus_dir = args.corpus_dir or Path(
            os.environ.get("SLIDE_TAGGER_CORPUS_DIR", "../mcp_slide_tagging/corpus")
        )
        _write_into_corpus(pptx, record, corpus_dir, fraction=args.fraction)

    if args.render:
        from slide_tagger.extractors.render import render_deck

        try:
            result = render_deck(pptx, poppler_path=args.poppler_path)
            print(
                f"# rendered {len(result.slides)} slide(s) -> {Path('data/renders') / result.deck_slug}",
                file=sys.stderr,
            )
        except (LibreOfficeNotFound, RenderError) as exc:
            print(f"# render skipped: {exc}", file=sys.stderr)
        except Exception as exc:  # poppler / pdf2image failures
            print(f"# render skipped (is poppler installed? --poppler-path): {exc}", file=sys.stderr)

    return 0


def _resolve_paste_deck(deck_arg: str, sources_dir: Path) -> tuple[Path, str] | int:
    """Resolve a `pack`/`ingest` positional `deck` (either a .pptx path or a stem
    that lives in `sources_dir`) → `(pptx_path, deck_slug)`. Mirrors `_resolve_deck`
    + `_BENCH_DEFAULT` stem lookup so the paste commands accept whichever form is
    handier."""
    arg_path = Path(deck_arg)
    if arg_path.suffix.lower() == ".pptx" and arg_path.exists():
        return arg_path, deck_slug(arg_path.name)
    # Treat as a stem.
    candidate = sources_dir / f"{deck_arg}.pptx"
    if not candidate.exists():
        print(
            f"Deck not found: tried {arg_path} (as .pptx) and {candidate} (as stem in {sources_dir}).",
            file=sys.stderr,
        )
        return 2
    return candidate, deck_slug(candidate.name)


def _cmd_pack(args: argparse.Namespace) -> int:
    """Build the paste bundle: write `in.md` + `meta.json` under
    `data/paste/<deck>/<variant>/`. The bundle is what the user pastes into
    claude.ai (or any other VLM web UI) to reproduce the API enrichment for $0."""
    from slide_tagger import paste as paste_mod

    resolved = _resolve_paste_deck(args.deck, args.sources)
    if isinstance(resolved, int):
        return resolved
    pptx, slug = resolved

    try:
        artifact = resolve_prompt(args.prompt)
    except FileNotFoundError as exc:
        print(str(exc), file=sys.stderr)
        return 2

    template = _build_template(pptx)
    in_path, meta = paste_mod.write_pack(
        artifact=artifact, template=template, deck_slug=slug,
        variant=args.variant, source_pptx=pptx, base=args.paste_dir,
    )
    print(f"# wrote paste bundle  -> {in_path}", file=sys.stderr)
    print(
        f"# prompt_version={artifact.version}  variant={args.variant}  "
        f"next_run={paste_mod.next_run_index(paste_mod.variant_dir(slug, args.variant, args.paste_dir))}",
        file=sys.stderr,
    )
    # Echo the path on stdout so callers can pipe / open it.
    print(in_path)
    return 0


def _cmd_ingest(args: argparse.Namespace) -> int:
    """Capture a VLM web-UI response, sanitize enums, re-impose Pipeline A fields,
    stamp provenance, and write `run_N.json` (auto-incrementing N)."""
    from slide_tagger import paste as paste_mod

    try:
        meta = paste_mod.load_meta(args.deck, args.variant, args.paste_dir)
    except FileNotFoundError as exc:
        print(str(exc), file=sys.stderr)
        return 2

    pptx = Path(meta.source_pptx)
    if not pptx.exists():
        print(
            f"# meta references missing source pptx: {pptx} (re-run `pack` after restoring it).",
            file=sys.stderr,
        )
        return 2

    try:
        artifact = resolve_prompt(args.prompt)
    except FileNotFoundError as exc:
        print(str(exc), file=sys.stderr)
        return 2
    # A prompt-version mismatch usually means the user edited the prompt between
    # pack and ingest — warn so a score-paste delta isn't silently mis-attributed.
    if artifact.version != meta.prompt_version:
        print(
            f"# warning: prompt_version changed since pack "
            f"({meta.prompt_version} -> {artifact.version}); record is stamped "
            f"with the CURRENT prompt version.",
            file=sys.stderr,
        )

    try:
        vlm_output = paste_mod.read_vlm_output(args.reply)
    except FileNotFoundError as exc:
        print(f"Reply not found: {exc}", file=sys.stderr)
        return 2
    except ValueError as exc:  # extract_json couldn't find a JSON object
        print(f"Could not parse VLM reply as JSON: {exc}", file=sys.stderr)
        return 1

    template = _build_template(pptx)
    template_core = {k: v for k, v in template.items() if k != "_legend"}

    out_path, record, changes = paste_mod.ingest_run(
        vlm_output=vlm_output, template_core=template_core, artifact=artifact,
        deck_slug=args.deck, variant=args.variant, model=args.model,
        tagged_by=args.tagged_by, base=args.paste_dir,
    )

    try:
        DeckTag.model_validate(record)
    except ValidationError as exc:
        print(
            f"# {len(exc.errors())} schema issue(s) after merge (not structural):",
            file=sys.stderr,
        )
        for err in exc.errors()[:10]:
            print(f"#   {'.'.join(str(x) for x in err['loc'])}: {err['msg']}", file=sys.stderr)
        if args.strict:
            return 1

    n_low = len(record["provenance"]["low_confidence_fields"])
    print(f"# wrote {out_path}", file=sys.stderr)
    print(
        f"# enum changes={len(changes)}  low_confidence_fields={n_low}  "
        f"prompt_version={artifact.version}",
        file=sys.stderr,
    )
    print(out_path)
    return 0


def _cmd_score_paste(args: argparse.Namespace) -> int:
    """Score every `run_N.json` for a variant against the hand-label answer key.
    Writes the scorecard to `data/paste/<deck>/<variant>/score.md` and echoes it
    to stdout. Mirrors `bench`'s per-deck reporting but on the paste runs."""
    from slide_tagger import paste as paste_mod

    runs = paste_mod.list_runs(args.deck, args.variant, args.paste_dir)
    if not runs:
        print(
            f"No runs in {paste_mod.variant_dir(args.deck, args.variant, args.paste_dir)}. "
            f"Run `slide-tagger ingest {args.deck} --variant {args.variant} <reply.json>` first.",
            file=sys.stderr,
        )
        return 2

    label_path = args.labels / f"{args.deck}.tagged.json"
    if not label_path.exists():
        print(f"Hand-label not found: {label_path}", file=sys.stderr)
        return 2
    truth = _load_tag(label_path)
    if isinstance(truth, int):
        return truth

    deck_scores = []
    for run in runs:
        pred = _load_tag(run)
        if isinstance(pred, int):
            return pred
        deck_scores.append(score_deck(pred, truth, name=run.stem))

    corpus = score_corpus(deck_scores)
    report = render_console(corpus)
    print(report)

    score_path = paste_mod.variant_dir(args.deck, args.variant, args.paste_dir) / "score.md"
    score_path.write_text(render_markdown(corpus), encoding="utf-8")
    print(f"\n# wrote {score_path}", file=sys.stderr)
    return 0


def _cmd_compare_paste(args: argparse.Namespace) -> int:
    """Diff per-field accuracy across two named variants for one deck. Tells you
    concretely 'v2's `message_type` is +6pp vs v1, `slide_purpose` is -2pp, …'."""
    from slide_tagger import paste as paste_mod

    variants = [v.strip() for v in args.variants.split(",") if v.strip()]
    if len(variants) != 2:
        print(
            f"--variants needs exactly two comma-separated names; got {variants}.",
            file=sys.stderr,
        )
        return 2

    label_path = args.labels / f"{args.deck}.tagged.json"
    if not label_path.exists():
        print(f"Hand-label not found: {label_path}", file=sys.stderr)
        return 2
    truth = _load_tag(label_path)
    if isinstance(truth, int):
        return truth

    def _variant_corpus(v: str):
        runs = paste_mod.list_runs(args.deck, v, args.paste_dir)
        if not runs:
            print(
                f"No runs for variant {v!r} (looked in "
                f"{paste_mod.variant_dir(args.deck, v, args.paste_dir)}).",
                file=sys.stderr,
            )
            return None
        deck_scores = []
        for run in runs:
            pred = _load_tag(run)
            if isinstance(pred, int):
                return None
            deck_scores.append(score_deck(pred, truth, name=run.stem))
        return score_corpus(deck_scores)

    a_name, b_name = variants
    a, b = _variant_corpus(a_name), _variant_corpus(b_name)
    if a is None or b is None:
        return 2

    all_fields = sorted(set(a.results) | set(b.results))

    def _n(corpus) -> int:
        return sum(r.scored for r in corpus.results.values())

    a_acc = a.headline_accuracy or 0.0
    b_acc = b.headline_accuracy or 0.0
    print(f"# compare-paste  deck={args.deck}  variants={a_name!r} vs {b_name!r}")
    print(f"# {a_name}: corpus headline {a_acc:.3f}  (n={_n(a)}, runs={len(a.deck_names)})")
    print(f"# {b_name}: corpus headline {b_acc:.3f}  (n={_n(b)}, runs={len(b.deck_names)})")
    print(f"# Δ headline: {b_acc - a_acc:+.3f}  ({b_name} − {a_name})")
    print()
    header = f"{'field':40s}  {a_name:>10s}  {b_name:>10s}  {'delta':>8s}"
    print(header)
    print("-" * len(header))
    for path in all_fields:
        ra, rb = a.results.get(path), b.results.get(path)
        aa = ra.accuracy if ra and ra.accuracy is not None else None
        ab = rb.accuracy if rb and rb.accuracy is not None else None
        d = (ab - aa) if (aa is not None and ab is not None) else None
        print(
            f"{path:40s}  "
            f"{('--' if aa is None else f'{aa:.3f}'):>10s}  "
            f"{('--' if ab is None else f'{ab:.3f}'):>10s}  "
            f"{('--' if d is None else f'{d:+.3f}'):>8s}"
        )
    return 0


def _cmd_score(args: argparse.Namespace) -> int:
    pred = _load_tag(args.predicted)
    if isinstance(pred, int):
        return pred
    truth = _load_tag(args.truth)
    if isinstance(truth, int):
        return truth
    deck = score_deck(pred, truth, name=args.predicted.stem)
    print(render_console(score_corpus([deck])))
    return 0


def _cmd_eval(args: argparse.Namespace) -> int:
    labels_dir: Path = args.labels
    pred_dir: Path = args.pred_dir
    if not labels_dir.is_dir():
        print(f"Labels dir not found: {labels_dir}", file=sys.stderr)
        return 2
    if not pred_dir.is_dir():
        print(f"Predictions dir not found: {pred_dir}", file=sys.stderr)
        return 2

    label_files = sorted(labels_dir.glob("*.json"))
    if not label_files:
        print(f"No hand-label JSON files in {labels_dir}", file=sys.stderr)
        return 2

    deck_scores = []
    for label_file in label_files:
        pred_file = pred_dir / label_file.name  # predictions share the hand-label filename
        if not pred_file.exists():
            print(f"# skip {label_file.name}: no prediction at {pred_file}", file=sys.stderr)
            continue
        truth = _load_tag(label_file)
        if isinstance(truth, int):
            return truth
        pred = _load_tag(pred_file)
        if isinstance(pred, int):
            return pred
        deck_scores.append(score_deck(pred, truth, name=label_file.stem))

    if not deck_scores:
        print("No matched (prediction, hand-label) pairs found.", file=sys.stderr)
        return 2

    corpus = score_corpus(deck_scores)
    print(render_console(corpus))
    if args.json:
        Path(args.json).write_text(render_json(corpus), encoding="utf-8")
        print(f"\n# wrote JSON report -> {args.json}", file=sys.stderr)
    if args.markdown:
        Path(args.markdown).write_text(render_markdown(corpus), encoding="utf-8")
        print(f"# wrote Markdown report -> {args.markdown}", file=sys.stderr)
    return 0


def main(argv: list[str] | None = None) -> int:
    # Emit UTF-8 regardless of the console code page (e.g. Windows cp932), so
    # the em-dash in the block header and non-ASCII slide titles print cleanly.
    for stream in (sys.stdout, sys.stderr):
        try:
            stream.reconfigure(encoding="utf-8")
        except (AttributeError, ValueError):
            pass

    # Load a local .env (e.g. ANTHROPIC_API_KEY for `bench`) if present. Real env
    # vars win over .env; missing python-dotenv is non-fatal.
    try:
        from dotenv import load_dotenv

        load_dotenv()
    except ImportError:
        pass

    parser = argparse.ArgumentParser(
        prog="slide-tagger",
        description="Pipeline A: deterministic structural extraction for slide decks.",
    )
    sub = parser.add_subparsers(dest="command", required=True)

    p_tag = sub.add_parser("tag", help="Extract structural data from a .pptx deck.")
    p_tag.add_argument("deck", type=Path, help="Path to a .pptx file")
    p_tag.add_argument(
        "--slide", type=int, default=None, help="Only this slide index (0-based)"
    )
    p_tag.add_argument(
        "--json",
        action="store_true",
        help="Emit the full structural JSON instead of paste-ready blocks",
    )
    p_tag.set_defaults(func=_cmd_tag)

    p_sum = sub.add_parser(
        "deck-summary",
        help="Emit a whole-deck DECK SUMMARY block for the deck-level VLM pass.",
    )
    p_sum.add_argument("deck", type=Path, help="Path to a .pptx file")
    p_sum.add_argument(
        "--json",
        action="store_true",
        help="Emit the deck-summary JSON instead of the paste-ready block",
    )
    p_sum.set_defaults(func=_cmd_deck_summary)

    p_tpl = sub.add_parser(
        "template",
        help="Emit a blank hand-tagging template (structural filled, semantics null).",
    )
    p_tpl.add_argument(
        "source", type=Path, help="A .pptx deck or a Pipeline A structural .json"
    )
    p_tpl.set_defaults(func=_cmd_template)

    p_val = sub.add_parser(
        "validate", help="Validate a hand-tagged JSON file and report completeness."
    )
    p_val.add_argument("labels", type=Path, help="Path to a hand-tagged .json file")
    p_val.set_defaults(func=_cmd_validate)

    p_render = sub.add_parser(
        "render",
        help="Render a .pptx to per-slide PNGs (full + thumbnail) via LibreOffice.",
    )
    p_render.add_argument("deck", type=Path, help="Path to a .pptx file")
    p_render.add_argument(
        "--out",
        type=Path,
        default=Path("data/renders"),
        help="Render output root (default: data/renders)",
    )
    p_render.add_argument(
        "--dpi", type=int, default=150, help="Full-render DPI (default: 150)"
    )
    p_render.add_argument(
        "--thumb", type=int, default=512, help="Thumbnail long-edge px (default: 512)"
    )
    p_render.add_argument(
        "--slide", type=int, default=None, help="Only this slide index (0-based)"
    )
    p_render.add_argument(
        "--poppler-path",
        dest="poppler_path",
        type=str,
        default=None,
        help="Path to poppler's bin/ if pdftoppm is not on PATH (common on Windows)",
    )
    p_render.set_defaults(func=_cmd_render)

    p_assets = sub.add_parser(
        "extract-assets",
        help="Extract recurring branding images (logos/watermarks) to PNGs and "
        "recurring_elements (scans slides + masters/layouts, recurses groups).",
    )
    p_assets.add_argument("deck", type=Path, help="Path to a .pptx file")
    p_assets.add_argument(
        "--out", type=Path, default=Path("reference_data/assets"),
        help="Asset output root (default: reference_data/assets)",
    )
    p_assets.add_argument(
        "--fraction", type=float, default=0.25,
        help="Min slide-coverage to treat a slide image as recurring (default: 0.25)",
    )
    p_assets.add_argument(
        "--into", type=Path, default=None,
        help="Merge the extracted elements into this hand-label JSON (dedupe by phash)",
    )
    p_assets.add_argument(
        "--contact-sheet", dest="contact_sheet", action="store_true",
        help="Also write a captioned grid (_contactsheet.png) of the extracted images",
    )
    p_assets.set_defaults(func=_cmd_extract_assets)

    p_merge = sub.add_parser(
        "merge",
        help="Re-impose Pipeline A structural fields from the template onto a VLM "
        "output (so the VLM can't corrupt deterministic fields).",
    )
    p_merge.add_argument(
        "vlm_output", type=Path, help="The enriched JSON returned by the VLM"
    )
    p_merge.add_argument(
        "template",
        type=Path,
        help="The Pipeline A template (input.json) the VLM enriched",
    )
    p_merge.add_argument(
        "-o",
        "--out",
        type=Path,
        default=None,
        help="Write merged JSON here (default: stdout). Status goes to stderr.",
    )
    p_merge.set_defaults(func=_cmd_merge)

    p_score = sub.add_parser(
        "score",
        help="Score one enriched (VLM) deck against its hand-label, field by field.",
    )
    p_score.add_argument("predicted", type=Path, help="Predicted (VLM-enriched) .json")
    p_score.add_argument("truth", type=Path, help="Hand-label ground-truth .json")
    p_score.set_defaults(func=_cmd_score)

    p_eval = sub.add_parser(
        "eval",
        help="Score every prediction against its hand-label and print a corpus scorecard.",
    )
    p_eval.add_argument(
        "--labels",
        type=Path,
        default=Path("reference_data/hand_labels"),
        help="Hand-label dir (default: reference_data/hand_labels)",
    )
    p_eval.add_argument(
        "--pred-dir",
        dest="pred_dir",
        type=Path,
        default=Path("data/tagged"),
        help="Predictions dir; files share the hand-label filename (default: data/tagged)",
    )
    p_eval.add_argument(
        "--json", type=Path, default=None, help="Also write a JSON report to this path"
    )
    p_eval.add_argument(
        "--markdown",
        type=Path,
        default=None,
        help="Also write a Markdown report to this path",
    )
    p_eval.set_defaults(func=_cmd_eval)

    p_bench = sub.add_parser(
        "bench",
        help="Run the enrichment prompt via the Claude API N times per deck and "
        "report mean ± std accuracy (averages out run-to-run variance).",
    )
    p_bench.add_argument(
        "--deck", action="append", default=None,
        help=f"Hand-label stem to run (repeatable). Default: {', '.join(_BENCH_DEFAULT)}",
    )
    p_bench.add_argument("--runs", type=int, default=3, help="Runs per deck (default: 3)")
    p_bench.add_argument(
        "--model", default="claude-opus-4-7",
        help="Model (default: claude-opus-4-7; pass claude-sonnet-4-6 for lower cost)",
    )
    p_bench.add_argument(
        "--effort", default="high", choices=["low", "medium", "high", "max"],
        help="Thinking/output effort (default: high)",
    )
    p_bench.add_argument(
        "--prompt", type=Path, default=Path("docs/deck_tagging_prompt.md"),
        help="Prompt markdown; its '## The Prompt' body is used as the system prompt",
    )
    p_bench.add_argument(
        "--sources", type=Path, default=Path("data/source"),
        help="Dir holding each deck's .pptx and .pdf (default: data/source)",
    )
    p_bench.add_argument(
        "--labels", type=Path, default=Path("reference_data/hand_labels"),
        help="Hand-label dir (the answer key)",
    )
    p_bench.add_argument(
        "--out", type=Path, default=Path("data/tagged/bench"),
        help="Where to write each run's merged JSON (default: data/tagged/bench)",
    )
    p_bench.add_argument(
        "--timeout", type=float, default=900.0,
        help="Per-request timeout in seconds (default: 900); a hang errors out instead of stalling",
    )
    p_bench.add_argument(
        "--quiet", action="store_true",
        help="Suppress the per-chunk streaming heartbeat",
    )
    p_bench.add_argument(
        "--label", default=None,
        help="Short label for this sweep (e.g. 'v3-placeholder') recorded in the results log",
    )
    p_bench.add_argument(
        "--json", type=Path, default=None,
        help="Also write a machine-readable JSON report to this path",
    )
    p_bench.add_argument(
        "--log", type=Path, default=Path("logs/bench_results.md"),
        help="Append a results row here (auto-created; default: logs/bench_results.md)",
    )
    p_bench.set_defaults(func=_cmd_bench)

    p_enrich = sub.add_parser(
        "enrich",
        help="Auto-tag an UNLABELED .pptx via the Claude API (Pipeline A + B) and "
        "write a finished tagged JSON — no hand-label needed.",
    )
    p_enrich.add_argument("deck", type=Path, help="Path to a .pptx file")
    p_enrich.add_argument(
        "--pdf", type=Path, default=None,
        help="Pre-converted deck PDF to upload (skips LibreOffice conversion)",
    )
    p_enrich.add_argument(
        "--model", default="claude-opus-4-7",
        help="Model (default: claude-opus-4-7; pass claude-sonnet-4-6 for lower cost)",
    )
    p_enrich.add_argument(
        "--effort", default="high", choices=["low", "medium", "high", "max"],
        help="Thinking/output effort (default: high)",
    )
    p_enrich.add_argument(
        "--prompt", type=Path, default=None,
        help="Prompt markdown (default: $SLIDE_TAGGER_PROMPT or docs/deck_tagging_prompt.md)",
    )
    p_enrich.add_argument(
        "--timeout", type=float, default=900.0,
        help="Per-request timeout in seconds (default: 900)",
    )
    p_enrich.add_argument("--quiet", action="store_true", help="Suppress the streaming heartbeat")
    p_enrich.add_argument(
        "--out", type=Path, default=None,
        help="Output JSON path (default: data/tagged/<slug>.tagged.json)",
    )
    p_enrich.add_argument(
        "--strict", action="store_true",
        help="Exit nonzero if the merged record fails schema validation (default: warn + write)",
    )
    p_enrich.add_argument(
        "--tagged-by", dest="tagged_by", default=None,
        help="Provenance tagged_by (default: auto:<model>)",
    )
    p_enrich.add_argument(
        "--into-corpus", dest="into_corpus", action="store_true",
        help="Also copy the result (+ extracted logo assets) into the served corpus dir",
    )
    p_enrich.add_argument(
        "--corpus-dir", dest="corpus_dir", type=Path, default=None,
        help="Corpus dir for --into-corpus (default: $SLIDE_TAGGER_CORPUS_DIR or ../mcp_slide_tagging/corpus)",
    )
    p_enrich.add_argument(
        "--fraction", type=float, default=0.25,
        help="Min slide-coverage to treat a slide image as recurring (default: 0.25)",
    )
    p_enrich.add_argument(
        "--render", action="store_true",
        help="Also render slide PNGs (off by default — images aren't served yet)",
    )
    p_enrich.add_argument(
        "--poppler-path", dest="poppler_path", type=str, default=None,
        help="Path to poppler's bin/ for --render (common on Windows)",
    )
    p_enrich.set_defaults(func=_cmd_enrich)

    # --- Web-UI paste harness (zero-API-cost iteration) ----------------------
    paste_dir_help = "Paste-harness output root (default: data/paste)"

    p_pack = sub.add_parser(
        "pack",
        help="Build a paste-ready bundle (system prompt + grounding + template) "
        "for a deck — paste into claude.ai to enrich for $0 (instead of `enrich`).",
    )
    p_pack.add_argument(
        "deck", type=str,
        help="A .pptx path OR a deck stem looked up in --sources (e.g. nigeria-...-v1)",
    )
    p_pack.add_argument(
        "--variant", required=True,
        help="Short experiment label (e.g. 'baseline', 'v2-deck-context'); creates "
        "data/paste/<deck>/<variant>/",
    )
    p_pack.add_argument(
        "--prompt", type=Path, default=None,
        help="Prompt markdown (default: $SLIDE_TAGGER_PROMPT or docs/deck_tagging_prompt.md)",
    )
    p_pack.add_argument(
        "--sources", type=Path, default=Path("data/source"),
        help="Dir where stems resolve to <stem>.pptx (default: data/source)",
    )
    p_pack.add_argument(
        "--paste-dir", dest="paste_dir", type=Path, default=Path("data/paste"),
        help=paste_dir_help,
    )
    p_pack.set_defaults(func=_cmd_pack)

    p_ingest = sub.add_parser(
        "ingest",
        help="Capture a VLM web-UI reply for a packed variant, sanitize + merge + "
        "stamp provenance, and write the next run_N.json under the variant dir.",
    )
    p_ingest.add_argument(
        "deck", type=str, help="Deck slug (as printed by `pack`)",
    )
    p_ingest.add_argument(
        "--variant", required=True, help="Same variant label used with `pack`",
    )
    p_ingest.add_argument(
        "reply", type=str,
        help="Path to a JSON file with the model's reply, or '-' to read from stdin",
    )
    p_ingest.add_argument(
        "--prompt", type=Path, default=None,
        help="Prompt markdown (default: $SLIDE_TAGGER_PROMPT or docs/deck_tagging_prompt.md)",
    )
    p_ingest.add_argument(
        "--model", default="claude-ai-web",
        help="Model name to stamp in provenance (default: claude-ai-web)",
    )
    p_ingest.add_argument(
        "--tagged-by", dest="tagged_by", default=None,
        help="Provenance tagged_by (default: paste:<model>:<variant>)",
    )
    p_ingest.add_argument(
        "--paste-dir", dest="paste_dir", type=Path, default=Path("data/paste"),
        help=paste_dir_help,
    )
    p_ingest.add_argument(
        "--strict", action="store_true",
        help="Exit nonzero if the merged record fails schema validation (default: warn + write)",
    )
    p_ingest.set_defaults(func=_cmd_ingest)

    p_scorep = sub.add_parser(
        "score-paste",
        help="Score every run_N.json for a paste variant against its hand-label; "
        "writes a Markdown scorecard alongside the runs.",
    )
    p_scorep.add_argument("deck", type=str, help="Deck slug")
    p_scorep.add_argument("--variant", required=True, help="Variant label")
    p_scorep.add_argument(
        "--labels", type=Path, default=Path("reference_data/hand_labels"),
        help="Hand-label dir (default: reference_data/hand_labels)",
    )
    p_scorep.add_argument(
        "--paste-dir", dest="paste_dir", type=Path, default=Path("data/paste"),
        help=paste_dir_help,
    )
    p_scorep.set_defaults(func=_cmd_score_paste)

    p_cmp = sub.add_parser(
        "compare-paste",
        help="Diff per-field accuracy across two paste variants for one deck "
        "(answers: did variant v2 actually help vs v1?).",
    )
    p_cmp.add_argument("deck", type=str, help="Deck slug")
    p_cmp.add_argument(
        "--variants", required=True,
        help="Two comma-separated variant names to compare, e.g. 'baseline,v2-deck-context'",
    )
    p_cmp.add_argument(
        "--labels", type=Path, default=Path("reference_data/hand_labels"),
        help="Hand-label dir (default: reference_data/hand_labels)",
    )
    p_cmp.add_argument(
        "--paste-dir", dest="paste_dir", type=Path, default=Path("data/paste"),
        help=paste_dir_help,
    )
    p_cmp.set_defaults(func=_cmd_compare_paste)

    args = parser.parse_args(argv)
    return args.func(args)


if __name__ == "__main__":
    sys.exit(main())
