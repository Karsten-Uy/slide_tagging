"""Deck-level design-system extraction: modal fonts/colors and pHash-based
recurring-element detection.

Deterministic (Pipeline A) — reads font/color data straight from the file and
groups repeated images by perceptual hash. The subjective parts (`grid`, and
*what* each recurring element is) are left for hand-labeling.
"""

from __future__ import annotations

import re
from collections import Counter
from dataclasses import dataclass
from io import BytesIO
from math import ceil

import imagehash
from PIL import Image
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN

from slide_tagger.schema.enums import (
    FontWeight,
    Position,
    TextAlignment,
)
from slide_tagger.schema.models import (
    ColorPalette,
    DesignSystem,
    RecurringElement,
    TextStyle,
)

_PICTURE_TYPES = {MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.LINKED_PICTURE}
_HASH_DISTANCE = 6  # max hamming distance to treat two images as "the same"
_RECURRING_FRACTION = 0.6  # appears on >=60% of slides -> recurring (init.md)

_ALIGN_MAP = {
    PP_ALIGN.LEFT: TextAlignment.LEFT,
    PP_ALIGN.CENTER: TextAlignment.CENTER,
    PP_ALIGN.RIGHT: TextAlignment.RIGHT,
    PP_ALIGN.JUSTIFY: TextAlignment.JUSTIFY,
}

_QUADRANTS = [
    [Position.TOP_LEFT, Position.TOP_CENTER, Position.TOP_RIGHT],
    [Position.MIDDLE_LEFT, Position.CENTER, Position.MIDDLE_RIGHT],
    [Position.BOTTOM_LEFT, Position.BOTTOM_CENTER, Position.BOTTOM_RIGHT],
]


@dataclass
class RawText:
    is_title: bool
    style: TextStyle


@dataclass
class RawImage:
    hash: imagehash.ImageHash
    position: Position | None
    slide_index: int


# --- low-level, best-effort field extraction (any of these can be unavailable) -


def _position(shape, slide_w: int, slide_h: int) -> Position | None:
    if shape.left is None or shape.top is None or not slide_w or not slide_h:
        return None
    cx = (shape.left + (shape.width or 0) / 2) / slide_w
    cy = (shape.top + (shape.height or 0) / 2) / slide_h
    col = 0 if cx < 1 / 3 else (1 if cx < 2 / 3 else 2)
    row = 0 if cy < 1 / 3 else (1 if cy < 2 / 3 else 2)
    return _QUADRANTS[row][col]


_PS_STYLE = re.compile(r"-(?:Bold|Italic|Oblique|Regular|Light|Medium|Semibold|Black|Thin)", re.I)
_FOUNDRY_TAIL = re.compile(r"\s*(?:MT|PS)$", re.I)  # "Arial MT" / "ArialMT"
_FOUNDRY_QUAL = re.compile(r"\s+(?:Pro|Std|LT|Com)\b.*$", re.I)  # "Helvetica Neue LT Std"


def _normalize_font(name: str | None) -> str | None:
    """Map a PostScript/foundry font name to its installable family so that
    `run.font.name = …` matches an installed font downstream (otherwise it silently
    falls back). 'Arial MT'/'ArialMT'/'Arial-BoldMT' -> 'Arial';
    'Helvetica Neue LT Std' -> 'Helvetica Neue'. 'Georgia'/'Calibri'/None unchanged."""
    if not name:
        return name
    n = _PS_STYLE.split(name.strip(), maxsplit=1)[0]
    n = _FOUNDRY_TAIL.sub("", n)
    n = _FOUNDRY_QUAL.sub("", n)
    return n.strip() or name.strip()


def _size_pt(font) -> float | None:
    try:
        return round(font.size.pt, 1) if font.size is not None else None
    except (AttributeError, ValueError):
        return None


def _font_hex(font) -> str | None:
    try:
        rgb = font.color.rgb
    except (AttributeError, TypeError):
        return None
    return f"#{rgb}" if rgb is not None else None


def _fill_hex(shape) -> str | None:
    try:
        return f"#{shape.fill.fore_color.rgb}"
    except (AttributeError, TypeError, ValueError, NotImplementedError):
        return None


def _align(alignment) -> TextAlignment | None:
    return _ALIGN_MAP.get(alignment)


def _text_style(shape) -> TextStyle | None:
    tf = shape.text_frame
    para = next((p for p in tf.paragraphs if p.text.strip()), None)
    if para is None:
        return None
    run = para.runs[0] if para.runs else None
    font = run.font if run is not None else para.font
    return TextStyle(
        font_family=_normalize_font(font.name),
        size_pt=_size_pt(font),
        weight=FontWeight.BOLD if font.bold else FontWeight.REGULAR,
        color_hex=_font_hex(font),
        alignment=_align(para.alignment),
    )


def _image_hash(shape) -> imagehash.ImageHash | None:
    try:
        with Image.open(BytesIO(shape.image.blob)) as im:
            return imagehash.phash(im.convert("RGB"))
    except Exception:  # EMF/WMF, corrupt, or non-raster image — skip
        return None


def extract_slide_design(
    slide, slide_w: int, slide_h: int, slide_index: int
) -> tuple[list[RawText], list[str], list[RawImage]]:
    """Collect raw font/color/image records for one slide."""
    texts: list[RawText] = []
    fills: list[str] = []
    images: list[RawImage] = []

    title_shape = slide.shapes.title
    # Compare shape_id, not object identity: python-pptx returns fresh wrapper
    # objects on each access, so `shape is title_shape` is always False.
    title_id = title_shape.shape_id if title_shape is not None else None
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip():
            style = _text_style(shape)
            if style is not None:
                is_title = title_id is not None and shape.shape_id == title_id
                texts.append(RawText(is_title=is_title, style=style))
        fill = _fill_hex(shape)
        if fill:
            fills.append(fill)
        if shape.shape_type in _PICTURE_TYPES:
            h = _image_hash(shape)
            if h is not None:
                images.append(
                    RawImage(hash=h, position=_position(shape, slide_w, slide_h),
                             slide_index=slide_index)
                )
    return texts, fills, images


# --- aggregation ------------------------------------------------------------


def _modal(values):
    vals = [v for v in values if v is not None]
    return Counter(vals).most_common(1)[0][0] if vals else None


def _modal_style(styles: list[TextStyle]) -> TextStyle:
    return TextStyle(
        font_family=_modal([s.font_family for s in styles]),
        size_pt=_modal([s.size_pt for s in styles]),
        weight=_modal([s.weight for s in styles]),
        color_hex=_modal([s.color_hex for s in styles]),
        alignment=_modal([s.alignment for s in styles]),
    )


def _rgb(hex_str: str) -> tuple[int, int, int]:
    h = hex_str.lstrip("#")
    return int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)


def _is_neutral(hex_str: str) -> bool:
    try:
        r, g, b = _rgb(hex_str)
    except (ValueError, IndexError):
        return False
    return max(r, g, b) - min(r, g, b) <= 16


def _palette(colors: list[str]) -> ColorPalette:
    counts = Counter(c for c in colors if c)
    neutrals: list[str] = []
    chromatic: list[str] = []
    for color, _ in counts.most_common():
        (neutrals if _is_neutral(color) else chromatic).append(color)
    primary = chromatic[0] if chromatic else (neutrals[0] if neutrals else None)
    accent = chromatic[1] if len(chromatic) > 1 else None
    return ColorPalette(primary=primary, accent=accent, neutrals=neutrals[:3])


def _recurring(images: list[RawImage], slide_count: int) -> list[RecurringElement]:
    groups: list[tuple[imagehash.ImageHash, list[RawImage]]] = []
    for im in images:
        for repr_hash, members in groups:
            if im.hash - repr_hash <= _HASH_DISTANCE:
                members.append(im)
                break
        else:
            groups.append((im.hash, [im]))

    threshold = max(2, ceil(_RECURRING_FRACTION * slide_count))
    out: list[RecurringElement] = []
    for repr_hash, members in groups:
        slides = sorted({m.slide_index for m in members})
        if len(slides) >= threshold:
            positions = [m.position for m in members if m.position]
            pos = Counter(positions).most_common(1)[0][0] if positions else None
            out.append(
                RecurringElement(
                    phash=str(repr_hash),
                    position=pos,
                    appears_on_slides=slides,
                    type=None,
                )
            )
    return out


def build_design_system(
    texts: list[RawText],
    fills: list[str],
    images: list[RawImage],
    slide_count: int,
) -> DesignSystem:
    """Aggregate raw per-slide records into a deck-level DesignSystem."""
    title_styles = [t.style for t in texts if t.is_title]
    body_styles = [t.style for t in texts if not t.is_title]
    text_colors = [t.style.color_hex for t in texts if t.style.color_hex]

    title_style = _modal_style(title_styles)
    body_style = _modal_style(body_styles)

    # Title placeholders are an unreliable font source: many decks put titles in
    # custom text boxes (title font ends up None) and the few that do carry one can
    # be a decorative outlier (e.g. "Georgia" on 3/39 runs while the deck is Arial).
    # So set the title *family* to the deck's dominant font (modal of body, then of
    # all text); keep the title's own size/weight/color/alignment. Body's modal is
    # reliable (lots of text) and only falls back if absent.
    dominant_font = body_style.font_family or _modal([t.style.font_family for t in texts])
    if dominant_font is not None:
        if title_style.font_family != dominant_font:
            title_style = title_style.model_copy(update={"font_family": dominant_font})
        if body_style.font_family is None:
            body_style = body_style.model_copy(update={"font_family": dominant_font})

    return DesignSystem(
        title_style=title_style,
        body_style=body_style,
        color_palette=_palette(text_colors + fills),
        default_text_alignment=_modal([t.style.alignment for t in texts]),
        grid=None,
        recurring_elements=_recurring(images, slide_count),
    )
