"""Pipeline A: deterministic structural extraction from a .pptx file.

No AI. Reads the file's own data with python-pptx, so the structural fields are
near-100% accurate. These are exactly the fields the VLM (Pipeline B) is told NOT
to guess — it receives them as grounding context instead.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from slide_tagger.extractors.structural.density import density_bucket
from slide_tagger.extractors.structural.design_system import (
    build_design_system,
    extract_slide_design,
)
from slide_tagger.schema.enums import Position, SourceFormat
from slide_tagger.schema.models import Density, DeckStructural, SlideStructural

# Fallback slide dimensions (EMU) if the deck doesn't declare them: 4:3 @ 10x7.5in.
_DEFAULT_W = 9144000
_DEFAULT_H = 6858000

_PICTURE_TYPES = {MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.LINKED_PICTURE}

_QUADRANTS = [
    [Position.TOP_LEFT, Position.TOP_CENTER, Position.TOP_RIGHT],
    [Position.MIDDLE_LEFT, Position.CENTER, Position.MIDDLE_RIGHT],
    [Position.BOTTOM_LEFT, Position.BOTTOM_CENTER, Position.BOTTOM_RIGHT],
]


def _quadrant(cx: float, cy: float) -> Position:
    col = 0 if cx < 1 / 3 else (1 if cx < 2 / 3 else 2)
    row = 0 if cy < 1 / 3 else (1 if cy < 2 / 3 else 2)
    return _QUADRANTS[row][col]


def _shape_area(shape) -> int:
    w = shape.width or 0
    h = shape.height or 0
    return int(w) * int(h)


def _shape_position(shape, slide_w: int, slide_h: int) -> Position | None:
    if shape.left is None or shape.top is None or not slide_w or not slide_h:
        return None
    cx = (shape.left + (shape.width or 0) / 2) / slide_w
    cy = (shape.top + (shape.height or 0) / 2) / slide_h
    return _quadrant(cx, cy)


def _has_text(shape) -> bool:
    return bool(getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip())


def _parse_slide(index: int, slide, slide_w: int, slide_h: int) -> SlideStructural:
    word_count = 0
    text_blocks = 0
    image_count = 0
    has_chart = False
    has_table = False
    content_area = 0

    title_text: str | None = None
    title_position: Position | None = None

    title_shape = slide.shapes.title  # title placeholder, or None
    if title_shape is not None and _has_text(title_shape):
        title_text = title_shape.text_frame.text.strip()
        title_position = _shape_position(title_shape, slide_w, slide_h)

    for shape in slide.shapes:
        is_content = False

        if _has_text(shape):
            word_count += len(shape.text_frame.text.split())
            text_blocks += 1
            is_content = True

        if shape.shape_type in _PICTURE_TYPES:
            image_count += 1
            is_content = True

        if getattr(shape, "has_chart", False):
            has_chart = True
            is_content = True

        if getattr(shape, "has_table", False):
            has_table = True
            is_content = True

        if is_content:
            content_area += _shape_area(shape)

    # visual_elements: distinct non-text content (images + a chart + a table).
    # Shapes/auto-shapes are not counted (often decorative); revisit if needed.
    visual_elements = image_count + (1 if has_chart else 0) + (1 if has_table else 0)

    slide_area = slide_w * slide_h
    coverage = content_area / slide_area if slide_area else 0.0
    whitespace = max(0.0, 1.0 - min(1.0, coverage))

    density = Density(
        word_count=word_count,
        text_blocks=text_blocks,
        visual_elements=visual_elements,
        whitespace_ratio_est=round(whitespace, 2),
        bucket=density_bucket(word_count, visual_elements),
    )

    return SlideStructural(
        index=index,
        title_text=title_text,
        title_position=title_position,
        image_count=image_count,
        has_chart=has_chart,
        has_table=has_table,
        density=density,
    )


def parse_pptx(path: str | Path) -> DeckStructural:
    """Parse a .pptx file into deterministic structural data."""
    path = Path(path)
    prs = Presentation(str(path))
    slide_w = prs.slide_width or _DEFAULT_W
    slide_h = prs.slide_height or _DEFAULT_H

    slides = []
    all_texts = []
    all_fills: list[str] = []
    all_images = []
    for i, slide in enumerate(prs.slides):
        slides.append(_parse_slide(i, slide, slide_w, slide_h))
        texts, fills, images = extract_slide_design(slide, slide_w, slide_h, i)
        all_texts.extend(texts)
        all_fills.extend(fills)
        all_images.extend(images)

    design_system = build_design_system(all_texts, all_fills, all_images, len(slides))

    return DeckStructural(
        source_filename=path.name,
        source_format=SourceFormat.PPTX,
        slide_count=len(slides),
        design_system=design_system,
        slides=slides,
    )
