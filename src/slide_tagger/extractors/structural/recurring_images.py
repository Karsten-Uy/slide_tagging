"""Detect and extract recurring *branding* images (logos, watermarks, recurring
graphics) from a .pptx, for the `extract-assets` command.

Builds on Pipeline A's image hashing (`design_system._image_hash` / `_position`)
but tuned for logos, which differ from generic recurring images three ways:
  1. they can be nested inside GROUP shapes (so we recurse into groups);
  2. they can live on the slide master / layouts (so we scan those too);
  3. they often recur on far fewer than the 60% that
     `design_system._RECURRING_FRACTION` requires (so we use a lower threshold).

Vector / grouped logos that have no raster blob can't be hashed or saved here —
those use the manual drop-in fallback (see docs/manual_tagging.md).
"""

from __future__ import annotations

from collections import Counter
from dataclasses import dataclass

import imagehash
from pptx.enum.shapes import MSO_SHAPE_TYPE

from slide_tagger.extractors.structural.design_system import (
    _HASH_DISTANCE,
    _PICTURE_TYPES,
    _image_hash,
    _position,
)
from slide_tagger.schema.enums import Position, RecurringElementType

_LOGO_RECURRING_FRACTION = 0.25  # logos recur on far fewer slides than 0.6
_LOGO_MAX_AREA_FRAC = 0.06  # a logo is small relative to the slide
_CORNERS = {
    Position.TOP_LEFT,
    Position.TOP_RIGHT,
    Position.BOTTOM_LEFT,
    Position.BOTTOM_RIGHT,
}


@dataclass
class RawLogoImage:
    phash: imagehash.ImageHash
    position: Position | None
    slide_index: int | None  # None for master/layout images
    source: str  # "slide" | "layout" | "master"
    blob: bytes
    area_frac: float


@dataclass
class LogoGroup:
    phash: str
    position: Position | None
    source: str
    slide_indices: list[int]
    area_frac: float
    blob: bytes
    type: RecurringElementType | None


def iter_picture_shapes(shapes):
    """Yield picture shapes from a shape collection, recursing into GROUP shapes."""
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_picture_shapes(shape.shapes)
        elif shape.shape_type in _PICTURE_TYPES:
            yield shape


def _blob(shape) -> bytes | None:
    try:
        return shape.image.blob
    except Exception:  # linked picture / no embedded blob
        return None


def _collect(shapes, *, source: str, slide_index: int | None, slide_w: int, slide_h: int) -> list[RawLogoImage]:
    out: list[RawLogoImage] = []
    area = (slide_w or 1) * (slide_h or 1)
    for shape in iter_picture_shapes(shapes):
        h = _image_hash(shape)  # None for EMF/WMF/corrupt (vector → manual fallback)
        if h is None:
            continue
        blob = _blob(shape)
        if blob is None:
            continue
        w, ht = (shape.width or 0), (shape.height or 0)
        out.append(
            RawLogoImage(
                phash=h,
                position=_position(shape, slide_w, slide_h),
                slide_index=slide_index,
                source=source,
                blob=blob,
                area_frac=(w * ht) / area if area else 0.0,
            )
        )
    return out


def collect_logo_images(prs) -> list[RawLogoImage]:
    """Raster pictures across slides + masters + layouts (groups recursed)."""
    sw, sh = prs.slide_width, prs.slide_height
    images: list[RawLogoImage] = []
    for i, slide in enumerate(prs.slides):
        images += _collect(slide.shapes, source="slide", slide_index=i, slide_w=sw, slide_h=sh)
    for master in prs.slide_masters:
        images += _collect(master.shapes, source="master", slide_index=None, slide_w=sw, slide_h=sh)
        for layout in master.slide_layouts:
            images += _collect(layout.shapes, source="layout", slide_index=None, slide_w=sw, slide_h=sh)
    return images


def auto_label(group: LogoGroup, slide_count: int, fraction: float) -> RecurringElementType | None:
    """Confidently label small corner branding as a logo; leave everything else to
    the human (never auto-guess watermark)."""
    recurs = len(group.slide_indices) / slide_count if slide_count else 0.0
    from_master = group.source in ("master", "layout")
    if (
        group.area_frac
        and group.area_frac < _LOGO_MAX_AREA_FRAC
        and group.position in _CORNERS
        and (from_master or recurs >= fraction)
    ):
        return RecurringElementType.LOGO
    return None


def cluster_recurring(
    images: list[RawLogoImage], slide_count: int, fraction: float = _LOGO_RECURRING_FRACTION
) -> list[LogoGroup]:
    """Greedy pHash clustering (hamming <= _HASH_DISTANCE). Keep a group when it
    recurs on >= `fraction` of slides OR is sourced from a master/layout (those
    render deck-wide by definition)."""
    reps: list[imagehash.ImageHash] = []
    groups: list[list[RawLogoImage]] = []
    for im in images:
        for rep, members in zip(reps, groups):
            if im.phash - rep <= _HASH_DISTANCE:
                members.append(im)
                break
        else:
            reps.append(im.phash)
            groups.append([im])

    out: list[LogoGroup] = []
    for members in groups:
        slide_idxs = sorted({m.slide_index for m in members if m.slide_index is not None})
        sources = {m.source for m in members}
        from_master = bool(sources & {"master", "layout"})
        coverage = len(slide_idxs) / slide_count if slide_count else 0.0
        if not (from_master or coverage >= fraction):
            continue
        positions = [m.position for m in members if m.position]
        rep = members[0]
        g = LogoGroup(
            phash=str(rep.phash),
            position=Counter(positions).most_common(1)[0][0] if positions else None,
            source="master" if "master" in sources else ("layout" if "layout" in sources else "slide"),
            slide_indices=slide_idxs,
            area_frac=rep.area_frac,
            blob=rep.blob,
            type=None,
        )
        g.type = auto_label(g, slide_count, fraction)
        out.append(g)
    return out


def extract_recurring_images(prs, fraction: float = _LOGO_RECURRING_FRACTION) -> list[LogoGroup]:
    """End-to-end: collect raster pictures (slides+masters+layouts, groups recursed),
    cluster by pHash, keep recurring/branding groups, and auto-label obvious logos."""
    return cluster_recurring(collect_logo_images(prs), len(prs.slides), fraction)
