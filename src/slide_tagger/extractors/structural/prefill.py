"""Pipeline A pre-fill: deterministically compute six per-slide fields that the
VLM was previously asked to derive.

Each inferrer is **conservative**: it returns a value only for unambiguous cases
and `None` (or omits the field) when judgment is needed — that way the VLM still
fills `Diagram` / `Icon-based` / `Framework graphic` for visuals, image-based
chart screenshots for `chart_type`, `subtitle` / `callout-box` / `citation` for
slot types, and `Bespoke` / `Broken` for `placeholder_compliance`.

Six fields covered (per the §A plan):

| Field                       | How                                                                   |
|-----------------------------|-----------------------------------------------------------------------|
| `dominant_visual_element`   | decision tree on image_count / has_chart / has_table / text presence  |
| `chart_type`                | `shape.chart.chart_type` (XL_CHART_TYPE) for native charts            |
| `embedded_data_present`     | True iff any shape has a native chart (vs a chart-screenshot image)   |
| `slot_types_present`        | shape inspection: title / chart / image / table / body / bullet /     |
|                             | footer (PP_PLACEHOLDER.FOOTER) / page_number (SLIDE_NUMBER) /         |
|                             | subtitle (SUBTITLE placeholder)                                       |
| `zones`                     | each content shape → `{name, region}` using a 5-band region mapping   |
|                             | from the existing `_shape_position` 3×3 quadrants                     |
| `placeholder_compliance`    | `shape.is_placeholder` ratio: ≥0.7 Pristine; 0.3-0.7 Reusable; else   |
|                             | left null (Bespoke/Broken stays VLM)                                  |

The module re-opens the `.pptx` (so callers don't have to thread a Presentation
object through), mirroring `extract_recurring_images`'s pattern. One extra parse
per deck is fine — this isn't a hot path.
"""

from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER

from slide_tagger.extractors.structural.pptx_parser import (
    _DEFAULT_H,
    _DEFAULT_W,
    _PICTURE_TYPES,
    _has_text,
    _shape_position,
)
from slide_tagger.schema.enums import (
    ChartType,
    DominantVisualElement,
    PlaceholderCompliance,
    Position,
    SlotType,
)


# --- region mapping: 3×3 quadrants → the 5-band names hand-labels use ---------

_REGION_BY_POSITION: dict[Position, str] = {
    Position.TOP_LEFT: "top-band",
    Position.TOP_CENTER: "top-band",
    Position.TOP_RIGHT: "top-band",
    Position.MIDDLE_LEFT: "left-main",
    Position.CENTER: "center",
    Position.MIDDLE_RIGHT: "right-callout",
    Position.BOTTOM_LEFT: "bottom-band",
    Position.BOTTOM_CENTER: "bottom-band",
    Position.BOTTOM_RIGHT: "bottom-band",
}


# --- chart_type: XL_CHART_TYPE → ChartType ------------------------------------

# Mapping is intentionally partial. Subtypes we don't recognize map to OTHER;
# native charts that python-pptx exposes as XL_CHART_TYPE values not handled
# here also fall through to OTHER (so the eval scorecard surfaces the gap
# instead of silently emitting null). The image-based chart case (no native
# chart on the slide despite a chart-looking image) is handled by leaving
# chart_type null at the slide level so the VLM still labels it.
_CHART_TYPE_MAP: dict[int, ChartType] = {
    # Line
    XL_CHART_TYPE.LINE: ChartType.LINE,
    XL_CHART_TYPE.LINE_MARKERS: ChartType.LINE,
    XL_CHART_TYPE.LINE_MARKERS_STACKED: ChartType.LINE,
    XL_CHART_TYPE.LINE_MARKERS_STACKED_100: ChartType.LINE,
    XL_CHART_TYPE.LINE_STACKED: ChartType.LINE,
    XL_CHART_TYPE.LINE_STACKED_100: ChartType.LINE,
    # Bar (single-series / clustered): include both column and bar orientations.
    XL_CHART_TYPE.BAR_CLUSTERED: ChartType.BAR,
    XL_CHART_TYPE.COLUMN_CLUSTERED: ChartType.BAR,
    XL_CHART_TYPE.BAR_OF_PIE: ChartType.BAR,
    XL_CHART_TYPE.THREE_D_BAR_CLUSTERED: ChartType.BAR,
    XL_CHART_TYPE.THREE_D_COLUMN_CLUSTERED: ChartType.BAR,
    # Stacked bar
    XL_CHART_TYPE.BAR_STACKED: ChartType.STACKED_BAR,
    XL_CHART_TYPE.BAR_STACKED_100: ChartType.STACKED_BAR,
    XL_CHART_TYPE.COLUMN_STACKED: ChartType.STACKED_BAR,
    XL_CHART_TYPE.COLUMN_STACKED_100: ChartType.STACKED_BAR,
    XL_CHART_TYPE.THREE_D_BAR_STACKED: ChartType.STACKED_BAR,
    XL_CHART_TYPE.THREE_D_BAR_STACKED_100: ChartType.STACKED_BAR,
    XL_CHART_TYPE.THREE_D_COLUMN_STACKED: ChartType.STACKED_BAR,
    XL_CHART_TYPE.THREE_D_COLUMN_STACKED_100: ChartType.STACKED_BAR,
    # Pie
    XL_CHART_TYPE.PIE: ChartType.PIE,
    XL_CHART_TYPE.PIE_EXPLODED: ChartType.PIE,
    XL_CHART_TYPE.PIE_OF_PIE: ChartType.PIE,
    XL_CHART_TYPE.THREE_D_PIE: ChartType.PIE,
    XL_CHART_TYPE.THREE_D_PIE_EXPLODED: ChartType.PIE,
    XL_CHART_TYPE.DOUGHNUT: ChartType.PIE,  # doughnut is a pie variant
    XL_CHART_TYPE.DOUGHNUT_EXPLODED: ChartType.PIE,
    # Scatter
    XL_CHART_TYPE.XY_SCATTER: ChartType.SCATTER,
    XL_CHART_TYPE.XY_SCATTER_LINES: ChartType.SCATTER,
    XL_CHART_TYPE.XY_SCATTER_LINES_NO_MARKERS: ChartType.SCATTER,
    XL_CHART_TYPE.XY_SCATTER_SMOOTH: ChartType.SCATTER,
    XL_CHART_TYPE.XY_SCATTER_SMOOTH_NO_MARKERS: ChartType.SCATTER,
    # Bubble
    XL_CHART_TYPE.BUBBLE: ChartType.BUBBLE,
    XL_CHART_TYPE.BUBBLE_THREE_D_EFFECT: ChartType.BUBBLE,
}


# --- per-slide inferrers ------------------------------------------------------


def _is_content_shape(shape) -> bool:
    """True iff the shape carries text, an image, a chart, or a table."""
    if _has_text(shape):
        return True
    if shape.shape_type in _PICTURE_TYPES:
        return True
    if getattr(shape, "has_chart", False):
        return True
    if getattr(shape, "has_table", False):
        return True
    return False


def _placeholder_type(shape):
    """Return `shape.placeholder_format.type` if the shape is a placeholder, else
    None. `placeholder_format` raises on non-placeholders in some python-pptx
    versions, so guard with `is_placeholder`."""
    if not getattr(shape, "is_placeholder", False):
        return None
    try:
        return shape.placeholder_format.type
    except (AttributeError, ValueError):
        return None


def infer_dominant_visual_element(slide) -> DominantVisualElement | None:
    """Return the dominant visual element when the file content is unambiguous;
    otherwise None (VLM picks Diagram / Icon-based / Framework graphic / Mixed)."""
    image_count = sum(1 for s in slide.shapes if s.shape_type in _PICTURE_TYPES)
    has_chart = any(getattr(s, "has_chart", False) for s in slide.shapes)
    has_table = any(getattr(s, "has_table", False) for s in slide.shapes)
    has_text = any(_has_text(s) for s in slide.shapes)

    # Pure single-type cases — high-confidence.
    if has_chart and not has_table and image_count == 0:
        return DominantVisualElement.CHART
    if has_table and not has_chart and image_count == 0:
        return DominantVisualElement.TABLE
    if image_count == 1 and not has_chart and not has_table:
        return DominantVisualElement.IMAGE
    if not has_chart and not has_table and image_count == 0 and has_text:
        return DominantVisualElement.PURE_TEXT
    # Multiple kinds of visual elements → MIXED (e.g. chart + image, table + image).
    distinct_kinds = sum([has_chart, has_table, image_count > 0])
    if distinct_kinds >= 2:
        return DominantVisualElement.MIXED
    # Ambiguous (multiple images and no other visual, no text, etc.) → defer.
    return None


def infer_chart_type(slide) -> ChartType | None:
    """Return the dominant **native** chart type on the slide. None when there's
    no native chart object — the VLM still needs to handle two cases we can't
    see from `python-pptx` alone: (a) chart screenshots (image of a chart, no
    native chart object → it'll write Bar/Line/etc), and (b) slides with no
    chart at all (it'll write N/A). Multiple native charts that disagree → also
    None (defer to VLM)."""
    charts = [s for s in slide.shapes if getattr(s, "has_chart", False)]
    if not charts:
        # We can't disambiguate "chart screenshot" from "no chart" without
        # vision, so leave null and let the VLM tag it.
        return None
    mapped = []
    for c in charts:
        try:
            ct = c.chart.chart_type
        except Exception:
            continue
        mapped.append(_CHART_TYPE_MAP.get(ct, ChartType.OTHER))
    if not mapped:
        return None
    if len(set(mapped)) == 1:
        return mapped[0]
    # Disagreement across native charts on the slide → don't guess.
    return None


def infer_embedded_data_present(slide) -> bool:
    """True iff the slide carries at least one **native** chart object (the chart
    data is in the .pptx, not just an image of a chart). False otherwise — even
    for slides with a screenshot of a chart, since the data isn't *embedded*."""
    return any(getattr(s, "has_chart", False) for s in slide.shapes)


def infer_slot_types_present(slide) -> list[SlotType]:
    """Detect the slot types Pipeline A can identify with high confidence. Skips
    `callout-box` and `citation` (need visual judgment); detects `subtitle`,
    `footer`, and `page-number` only via PP_PLACEHOLDER (text-heuristics aren't
    reliable enough)."""
    slots: set[SlotType] = set()

    # Title placeholder.
    title_shape = slide.shapes.title
    if title_shape is not None and _has_text(title_shape):
        slots.add(SlotType.TITLE)

    for shape in slide.shapes:
        if shape is title_shape:
            continue
        # Native visuals.
        if getattr(shape, "has_chart", False):
            slots.add(SlotType.CHART)
        if shape.shape_type in _PICTURE_TYPES:
            slots.add(SlotType.IMAGE)
        if getattr(shape, "has_table", False):
            slots.add(SlotType.TABLE)
        # Placeholder-driven slots (subtitle / footer / page-number).
        ph_type = _placeholder_type(shape)
        if ph_type == PP_PLACEHOLDER.SUBTITLE:
            slots.add(SlotType.SUBTITLE)
        elif ph_type == PP_PLACEHOLDER.FOOTER:
            slots.add(SlotType.FOOTER)
        elif ph_type == PP_PLACEHOLDER.SLIDE_NUMBER:
            slots.add(SlotType.PAGE_NUMBER)
        # Text shapes → body-text or bullet-list.
        if _has_text(shape):
            try:
                paragraphs = list(shape.text_frame.paragraphs)
            except Exception:
                paragraphs = []
            has_bullets = any(
                (p.level or 0) > 0 and p.text.strip() for p in paragraphs
            )
            has_body = any(p.text.strip() for p in paragraphs)
            if has_bullets:
                slots.add(SlotType.BULLET_LIST)
            elif has_body and ph_type not in (
                PP_PLACEHOLDER.SUBTITLE,
                PP_PLACEHOLDER.FOOTER,
                PP_PLACEHOLDER.SLIDE_NUMBER,
            ):
                slots.add(SlotType.BODY_TEXT)

    # Stable order so equal-set runs serialize identically (eval ENUM_LIST cares
    # about set equality, not order, but a stable order keeps diffs clean).
    return sorted(slots, key=lambda s: s.value)


def _zone_name_for(shape, title_shape) -> str:
    """Best-effort semantic name for a shape's zone — mirrors the hand-label
    vocabulary (title / subtitle / chart / table / image / body / footer /
    page-number / shape)."""
    ph_type = _placeholder_type(shape)
    # Title check: identity match OR title-placeholder type (python-pptx's
    # `slide.shapes.title` can return an object that isn't `is`-equal to the
    # one yielded by iterating `slide.shapes`).
    if shape is title_shape or ph_type in (
        PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE,
    ):
        return "title"
    if getattr(shape, "has_chart", False):
        return "chart"
    if getattr(shape, "has_table", False):
        return "table"
    if shape.shape_type in _PICTURE_TYPES:
        return "image"
    if ph_type == PP_PLACEHOLDER.SUBTITLE:
        return "subtitle"
    if ph_type == PP_PLACEHOLDER.FOOTER:
        return "footer"
    if ph_type == PP_PLACEHOLDER.SLIDE_NUMBER:
        return "page-number"
    if _has_text(shape):
        return "body"
    return "shape"


def infer_zones(slide, slide_w: int, slide_h: int) -> list[dict[str, str]]:
    """One Zone per content-bearing shape: {name, region}. Region comes from
    the existing 3×3 `_shape_position` mapped down to the 5 hand-label bands
    (top-band / left-main / center / right-callout / bottom-band).

    Zones for empty/decorative shapes are skipped (matches what hand-labels
    record). Duplicates within a slide are de-duped so we don't emit twenty
    {name=body, region=top-band} rows on a dense slide."""
    title_shape = slide.shapes.title
    zones: list[dict[str, str]] = []
    seen: set[tuple[str, str]] = set()
    for shape in slide.shapes:
        if not _is_content_shape(shape):
            continue
        pos = _shape_position(shape, slide_w, slide_h)
        if pos is None:
            continue
        region = _REGION_BY_POSITION.get(pos, pos.value)
        name = _zone_name_for(shape, title_shape)
        key = (name, region)
        if key in seen:
            continue
        seen.add(key)
        zones.append({"name": name, "region": region})
    return zones


def infer_placeholder_compliance(slide) -> PlaceholderCompliance | None:
    """Pristine (≥70% of content shapes come from the layout/master) or Reusable
    (30-70%). Below 30% → leave null so the VLM picks Bespoke vs Broken."""
    content = [s for s in slide.shapes if _is_content_shape(s)]
    if not content:
        return None
    placeholder_count = sum(1 for s in content if getattr(s, "is_placeholder", False))
    ratio = placeholder_count / len(content)
    if ratio >= 0.7:
        return PlaceholderCompliance.PRISTINE
    if ratio >= 0.3:
        return PlaceholderCompliance.REUSABLE
    return None


# --- deck-level entry point ---------------------------------------------------


def prefill_slide(slide, slide_w: int, slide_h: int) -> dict[str, Any]:
    """Compute the 6 pre-fillable fields for one python-pptx slide. Returns a
    dict using the wire-format strings (enum `.value`) so it can be merged
    straight into a SlideTag's JSON form."""
    out: dict[str, Any] = {}

    dve = infer_dominant_visual_element(slide)
    if dve is not None:
        out["dominant_visual_element"] = dve.value

    ct = infer_chart_type(slide)
    if ct is not None:
        out["chart_type"] = ct.value

    out["embedded_data_present"] = infer_embedded_data_present(slide)

    slots = infer_slot_types_present(slide)
    if slots:
        out["slot_types_present"] = [s.value for s in slots]

    zones = infer_zones(slide, slide_w, slide_h)
    if zones:
        out["zones"] = zones

    pc = infer_placeholder_compliance(slide)
    if pc is not None:
        out["placeholder_compliance"] = pc.value

    return out


def prefill_deck(path: str | Path) -> list[dict[str, Any]]:
    """Compute the 6 pre-fillable fields for every slide in a .pptx, returned
    in slide index order. One dict per slide (possibly empty when nothing was
    determinable). Re-opens the .pptx — small cost, keeps the API simple."""
    path = Path(path)
    prs = Presentation(str(path))
    slide_w = prs.slide_width or _DEFAULT_W
    slide_h = prs.slide_height or _DEFAULT_H
    return [prefill_slide(slide, slide_w, slide_h) for slide in prs.slides]
