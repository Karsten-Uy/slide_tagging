"""Tests for the §A Pipeline-A pre-fill — the 6 per-slide fields the VLM used
to derive (dominant_visual_element / chart_type / embedded_data_present /
slot_types_present / zones / placeholder_compliance).

Two layers:
- **Inferrer unit tests** against the small make_sample_deck fixture, asserting
  the deterministic value for each slide where Pipeline A is confident.
- **Pipeline integration**: _build_template populates the prefill fields, and
  merge_structural's `_SLIDE_PREFILL` guard re-imposes them when the VLM tries
  to overwrite a filled value (but yields to the VLM when Pipeline A left null).
"""

from __future__ import annotations

import importlib.util
import sys
from pathlib import Path

import pytest
from pptx import Presentation

from slide_tagger.cli import _build_template
from slide_tagger.extractors.structural.prefill import (
    _CHART_TYPE_MAP,
    infer_chart_type,
    infer_dominant_visual_element,
    infer_embedded_data_present,
    infer_placeholder_compliance,
    infer_slot_types_present,
    infer_zones,
    prefill_deck,
)
from slide_tagger.merge import _SLIDE_PREFILL, merge_structural
from slide_tagger.schema.enums import (
    ChartType,
    DominantVisualElement,
    PlaceholderCompliance,
    SlotType,
)

_ROOT = Path(__file__).resolve().parents[1]


def _load_sample_builder():
    spec = importlib.util.spec_from_file_location(
        "make_sample_deck", _ROOT / "scripts" / "make_sample_deck.py"
    )
    module = importlib.util.module_from_spec(spec)
    sys.modules["make_sample_deck"] = module
    spec.loader.exec_module(module)
    return module


@pytest.fixture(scope="module")
def sample_pptx(tmp_path_factory):
    builder = _load_sample_builder()
    path = tmp_path_factory.mktemp("prefilldeck") / "sample.pptx"
    builder.build(path)
    return path


@pytest.fixture(scope="module")
def sample_prs(sample_pptx):
    """Open the .pptx once for inferrer-level tests."""
    return Presentation(str(sample_pptx))


# --- inferrer unit tests --------------------------------------------------------
# The fixture has 4 slides: 0 Title, 1 Bulleted content, 2 Chart slide, 3 Section
# divider. Each one exercises a different prefill path.


def test_dominant_visual_element_chart_only(sample_prs):
    chart_slide = sample_prs.slides[2]  # title + native chart, no images/tables
    assert infer_dominant_visual_element(chart_slide) == DominantVisualElement.CHART


def test_dominant_visual_element_pure_text(sample_prs):
    # Slide 1 = bulleted content, just text — no chart/table/image.
    assert (
        infer_dominant_visual_element(sample_prs.slides[1])
        == DominantVisualElement.PURE_TEXT
    )


def test_chart_type_maps_column_clustered_to_bar(sample_prs):
    chart_slide = sample_prs.slides[2]
    assert infer_chart_type(chart_slide) == ChartType.BAR


def test_chart_type_none_when_no_native_chart(sample_prs):
    # The bulleted-content slide has no native chart object → prefill returns
    # None so the VLM can tag chart-screenshot images visually (or write N/A
    # for slides with no chart at all). Re-deriving chart-screenshot type from
    # the .pptx alone isn't possible.
    assert infer_chart_type(sample_prs.slides[1]) is None


def test_chart_type_map_covers_common_xl_chart_types():
    """Sanity check the mapping: every supported XL_CHART_TYPE resolves to a
    ChartType enum value (no stale ints, no typos)."""
    for chart_type, mapped in _CHART_TYPE_MAP.items():
        assert isinstance(mapped, ChartType), (chart_type, mapped)


def test_embedded_data_present_true_for_native_chart(sample_prs):
    assert infer_embedded_data_present(sample_prs.slides[2]) is True


def test_embedded_data_present_false_when_no_chart(sample_prs):
    assert infer_embedded_data_present(sample_prs.slides[1]) is False


def test_slot_types_title_only_for_title_slide(sample_prs):
    # Slide 0 = Title slide layout (TITLE + SUBTITLE placeholders).
    slots = infer_slot_types_present(sample_prs.slides[0])
    assert SlotType.TITLE in slots
    assert SlotType.SUBTITLE in slots
    # No body-text on a pure title slide.
    assert SlotType.CHART not in slots
    assert SlotType.IMAGE not in slots


def test_slot_types_detects_bullets_and_chart(sample_prs):
    bulleted = infer_slot_types_present(sample_prs.slides[1])
    assert SlotType.TITLE in bulleted
    # The bulleted-content slide has paragraphs at level 0 (no indent), so the
    # bullet detector should NOT classify it as bullet-list (python-pptx reports
    # level=0 for the default placeholder bullets). It should still be body-text.
    # If your sample changes to use level>0 explicitly, BULLET_LIST is fine too.
    assert SlotType.BODY_TEXT in bulleted or SlotType.BULLET_LIST in bulleted

    chart_slots = infer_slot_types_present(sample_prs.slides[2])
    assert SlotType.TITLE in chart_slots
    assert SlotType.CHART in chart_slots


def test_zones_dedupe_and_include_chart(sample_prs):
    slide_w = sample_prs.slide_width
    slide_h = sample_prs.slide_height
    zones = infer_zones(sample_prs.slides[2], slide_w, slide_h)
    names = {z["name"] for z in zones}
    assert "title" in names
    assert "chart" in names
    # Regions come from the 5-band mapping, not raw 3×3 quadrants.
    valid_regions = {"top-band", "left-main", "center", "right-callout", "bottom-band"}
    assert {z["region"] for z in zones} <= valid_regions
    # De-duped: no two zones with the same (name, region).
    pairs = [(z["name"], z["region"]) for z in zones]
    assert len(pairs) == len(set(pairs))


def test_placeholder_compliance_pristine_when_layout_placeholders(sample_prs):
    # Every slide in the fixture is built from python-pptx's standard layouts and
    # uses their placeholders (TITLE / SUBTITLE / BODY). The chart slide adds one
    # non-placeholder chart shape (still mostly placeholders).
    for s in sample_prs.slides:
        pc = infer_placeholder_compliance(s)
        # Either Pristine (all placeholders) or Reusable (chart slide). Never None
        # for these layout-driven slides.
        assert pc in {PlaceholderCompliance.PRISTINE, PlaceholderCompliance.REUSABLE}


def test_prefill_deck_returns_one_dict_per_slide(sample_pptx, sample_prs):
    out = prefill_deck(sample_pptx)
    assert len(out) == len(list(sample_prs.slides))
    # Every entry is a dict (possibly empty); embedded_data_present is always set.
    for entry in out:
        assert isinstance(entry, dict)
        assert "embedded_data_present" in entry


# --- pipeline integration: _build_template + merge guard -----------------------


def test_build_template_populates_prefill_fields(sample_pptx):
    tpl = _build_template(sample_pptx)
    chart_slide = next(s for s in tpl["slides"] if s["index"] == 2)
    assert chart_slide["dominant_visual_element"] == DominantVisualElement.CHART.value
    assert chart_slide["chart_type"] == ChartType.BAR.value
    assert chart_slide["embedded_data_present"] is True
    assert SlotType.CHART.value in chart_slide["slot_types_present"]
    assert any(z["name"] == "chart" for z in chart_slide["zones"])
    assert chart_slide["placeholder_compliance"] in {
        PlaceholderCompliance.PRISTINE.value,
        PlaceholderCompliance.REUSABLE.value,
    }


def test_merge_reimposes_prefill_when_template_has_value(sample_pptx):
    """If Pipeline A pre-filled a field, the VLM cannot overwrite it."""
    tpl = _build_template(sample_pptx)
    tpl_core = {k: v for k, v in tpl.items() if k != "_legend"}

    # Pretend the VLM returned a wrong dominant_visual_element on slide 2.
    vlm_out = {
        "source_filename": tpl_core["source_filename"],
        "source_format": tpl_core["source_format"],
        "slide_count": tpl_core["slide_count"],
        "slides": [
            {"index": s["index"], "dominant_visual_element": "Diagram"}
            for s in tpl_core["slides"]
        ],
    }
    merged = merge_structural(vlm_out, tpl_core)
    chart_merged = next(s for s in merged["slides"] if s["index"] == 2)
    # Pipeline A's CHART value wins over the VLM's "Diagram" — merge guard re-imposed.
    assert chart_merged["dominant_visual_element"] == DominantVisualElement.CHART.value


def test_merge_keeps_vlm_value_when_template_is_null(sample_pptx):
    """If Pipeline A left a field null (conservative), the VLM's answer is kept."""
    tpl = _build_template(sample_pptx)
    tpl_core = {k: v for k, v in tpl.items() if k != "_legend"}

    # Force the template's dominant_visual_element to null on slide 2 — simulates
    # an ambiguous case where prefill returned None.
    for s in tpl_core["slides"]:
        if s["index"] == 2:
            s["dominant_visual_element"] = None

    vlm_out = {
        "source_filename": tpl_core["source_filename"],
        "source_format": tpl_core["source_format"],
        "slide_count": tpl_core["slide_count"],
        "slides": [
            {"index": s["index"], "dominant_visual_element": "Diagram"}
            for s in tpl_core["slides"]
        ],
    }
    merged = merge_structural(vlm_out, tpl_core)
    chart_merged = next(s for s in merged["slides"] if s["index"] == 2)
    # Template was null → VLM's "Diagram" wins.
    assert chart_merged["dominant_visual_element"] == "Diagram"


def test_slide_prefill_set_matches_documented_fields():
    """The merge guard's _SLIDE_PREFILL must list the same 6 fields as the
    prompt section + prefill module so they stay in sync."""
    assert set(_SLIDE_PREFILL) == {
        "dominant_visual_element",
        "chart_type",
        "embedded_data_present",
        "slot_types_present",
        "zones",
        "placeholder_compliance",
    }
