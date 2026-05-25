"""Tests for deck-level design_system extraction (modal fonts/colors + pHash
recurring-element detection). Uses a dedicated fixture deck with a repeated logo
and styled titles, so it doesn't disturb the shared sample deck."""

from __future__ import annotations

from io import BytesIO

import pytest
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from slide_tagger.extractors.structural.design_system import (
    _normalize_font,
    build_design_system,
)
from slide_tagger.extractors.structural.design_system import RawText
from slide_tagger.extractors.structural.pptx_parser import parse_pptx
from slide_tagger.schema.enums import FontWeight
from slide_tagger.schema.models import TextStyle


@pytest.mark.parametrize(
    "raw, expected",
    [
        ("Arial MT", "Arial"),
        ("ArialMT", "Arial"),
        ("Arial-BoldMT", "Arial"),
        ("Helvetica Neue LT Std", "Helvetica Neue"),
        ("Georgia", "Georgia"),
        ("Calibri", "Calibri"),
        (None, None),
        ("", ""),
    ],
)
def test_normalize_font(raw, expected):
    assert _normalize_font(raw) == expected


def test_title_font_falls_back_to_body():
    # title placeholder has no font (common); body text is Arial.
    texts = [
        RawText(is_title=True, style=TextStyle(font_family=None, size_pt=32.0, weight=FontWeight.BOLD)),
        RawText(is_title=False, style=TextStyle(font_family="Arial", size_pt=12.0)),
    ]
    ds = build_design_system(texts, fills=[], images=[], slide_count=1)
    assert ds.title_style.font_family == "Arial"  # fell back to body font
    assert ds.title_style.size_pt == 32.0  # title's own size kept, not body's


def test_title_font_outlier_overridden_by_dominant():
    # a decorative outlier title font is overridden by the deck's dominant (body) font.
    texts = [
        RawText(is_title=True, style=TextStyle(font_family="Georgia", size_pt=32.0)),
        RawText(is_title=False, style=TextStyle(font_family="Arial", size_pt=12.0)),
        RawText(is_title=False, style=TextStyle(font_family="Arial", size_pt=12.0)),
    ]
    ds = build_design_system(texts, fills=[], images=[], slide_count=1)
    assert ds.title_style.font_family == "Arial"
    assert ds.title_style.size_pt == 32.0  # size preserved


def _logo_bytes() -> bytes:
    buf = BytesIO()
    Image.new("RGB", (120, 60), (200, 30, 30)).save(buf, format="PNG")
    return buf.getvalue()


@pytest.fixture(scope="module")
def deck(tmp_path_factory):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    logo = _logo_bytes()

    # 5 content slides; logo on 4 of them (top-right) -> recurring (>=60%).
    for i in range(5):
        s = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
        title = s.shapes.title
        title.text = f"Section {i}"
        run = title.text_frame.paragraphs[0].runs[0]
        run.font.size = Pt(32)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
        if i < 4:
            s.shapes.add_picture(
                BytesIO(logo), Inches(11.8), Inches(0.2), Inches(1.2), Inches(0.6)
            )

    path = tmp_path_factory.mktemp("ds") / "styled.pptx"
    prs.save(str(path))
    return parse_pptx(path)


def test_design_system_present(deck):
    assert deck.design_system is not None


def test_title_style_extracted(deck):
    ts = deck.design_system.title_style
    assert ts.size_pt == 32.0
    assert ts.weight is not None and ts.weight.value == "bold"
    assert ts.color_hex == "#1A1A1A"


def test_palette_has_colors(deck):
    palette = deck.design_system.color_palette
    # the dark title color should appear somewhere in the palette
    found = palette.primary in {"#1A1A1A"} or "#1A1A1A" in palette.neutrals
    assert found


def test_recurring_logo_detected(deck):
    recurring = deck.design_system.recurring_elements
    assert len(recurring) == 1
    el = recurring[0]
    assert el.appears_on_slides == [0, 1, 2, 3]
    assert el.position is not None
    assert el.type is None  # left for hand-labeling
    assert isinstance(el.phash, str) and el.phash


def test_non_recurring_not_flagged(deck):
    # only one repeated image group meets the 60% threshold
    assert all(len(r.appears_on_slides) >= 3 for r in deck.design_system.recurring_elements)
