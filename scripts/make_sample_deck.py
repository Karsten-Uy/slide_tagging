"""Generate a small, varied sample .pptx so Pipeline A can be run and tested
without needing a real deck. Slides cover a few distinct density profiles:

  0  title slide      -> sparse, no visuals
  1  bulleted content -> balanced, several text blocks
  2  data slide       -> has a chart (visual_elements >= 1)
  3  section divider  -> sparse, one short label

Usage:  python scripts/make_sample_deck.py [output.pptx]
"""

from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches


def build(path: Path) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)

    # 0: Title slide (layout 0 = Title Slide)
    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = "Q3 Business Review"
    s.placeholders[1].text = "Acme Corp · October 2026"

    # 1: Bulleted content (layout 1 = Title and Content)
    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = "Priorities for next quarter"
    body = s.placeholders[1].text_frame
    body.text = "Grow the enterprise pipeline"
    for line in (
        "Ship the new onboarding flow",
        "Reduce churn in the SMB segment",
        "Hire two senior engineers",
        "Launch in the EU region",
    ):
        body.add_paragraph().text = line

    # 2: Data slide with a chart (layout 5 = Title Only)
    s = prs.slides.add_slide(prs.slide_layouts[5])
    s.shapes.title.text = "Enterprise revenue grew 40% YoY"
    chart_data = CategoryChartData()
    chart_data.categories = ["Q1", "Q2", "Q3"]
    chart_data.add_series("Revenue ($M)", (3.2, 4.1, 5.7))
    s.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1), Inches(1.8), Inches(11), Inches(5),
        chart_data,
    )

    # 3: Section divider (layout 2 = Section Header)
    s = prs.slides.add_slide(prs.slide_layouts[2])
    s.shapes.title.text = "Appendix"

    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(path))
    return path


if __name__ == "__main__":
    out = Path(sys.argv[1]) if len(sys.argv) > 1 else Path("data/source/sample_deck.pptx")
    written = build(out)
    print(f"Wrote {written}")
